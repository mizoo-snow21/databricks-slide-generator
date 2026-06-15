"""Tests for PPTX slide generation service."""

from __future__ import annotations

import base64
from pathlib import Path
from unittest.mock import patch

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx.util import Inches

from services import pptx_slides_service
from services.html_slides_service import _merge_brand
from services.pptx_slides_service import (
    WIDGET_REF_KEYS,
    _NAMED_LAYOUT_PATTERNS,
    _blank_slide_layout,
    _build_slide,
    _fill_corporate_named_layout,
    _resolve_named_layout,
    generate_pptx_slides,
)

# Minimal valid 1x1 PNG for embedding / mocks
MINI_PNG = base64.standard_b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _slide_texts(prs: Presentation) -> list[str]:
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                t = (p.text or "").strip()
                if t:
                    texts.append(t)
    return texts


def _all_text_blob(prs: Presentation) -> str:
    return "\n".join(_slide_texts(prs))


def test_generate_pptx_returns_valid_file() -> None:
    path = generate_pptx_slides(
        "Deck",
        [{"layout": "title", "title": "Hello", "subtitle": "World"}],
    )
    assert path.exists()
    assert path.suffix == ".pptx"
    assert "genie-slide-" in path.name
    prs = Presentation(str(path))
    assert len(prs.slides) == 1


def test_generate_pptx_with_corporate_template_strips_existing_slides() -> None:
    tpl = (
        Path(__file__).resolve().parent.parent
        / "assets"
        / "databricks-corp-template.pptx"
    )
    if not tpl.is_file():
        pytest.skip("databricks-corp-template.pptx not in repo assets/")
    path = generate_pptx_slides(
        "Test",
        [{"layout": "title", "title": "Test", "subtitle": "sub"}],
        pptx_template_path=tpl,
    )
    assert path.exists()
    prs = Presentation(str(path))
    assert len(prs.slides) == 1
    assert len(prs.slide_masters) >= 1


def test_title_layout() -> None:
    path = generate_pptx_slides(
        "T",
        [{"layout": "title", "title": "Main Title", "subtitle": "Sub here"}],
    )
    prs = Presentation(str(path))
    blob = _all_text_blob(prs)
    assert "Main Title" in blob
    assert "Sub here" in blob
    # Title slide should have a background shape (gradient/solid fill)
    slide = prs.slides[0]
    assert len(slide.shapes) >= 1


def test_content_basic_layout() -> None:
    path = generate_pptx_slides(
        "C",
        [
            {
                "layout": "content_basic",
                "title": "Topic",
                "body": "Body line",
                "bullets": "a\nb;c",
            },
        ],
    )
    prs = Presentation(str(path))
    blob = _all_text_blob(prs)
    assert "Topic" in blob
    assert "Body line" in blob
    assert "a" in blob
    assert "b" in blob
    assert "c" in blob


def test_content_2col_layout() -> None:
    path = generate_pptx_slides(
        "2c",
        [
            {
                "layout": "content_2col",
                "title": "Two Cols",
                "col1_header": "H1",
                "col1_body": "B1",
                "col2_header": "H2",
                "col2_body": "B2",
            },
        ],
    )
    prs = Presentation(str(path))
    blob = _all_text_blob(prs)
    assert "Two Cols" in blob
    assert "H1" in blob and "B1" in blob
    assert "H2" in blob and "B2" in blob


def test_section_break_layout() -> None:
    path = generate_pptx_slides(
        "S",
        [{"layout": "section_break", "title": "Section A"}],
    )
    prs = Presentation(str(path))
    assert "Section A" in _all_text_blob(prs)


def test_chart_embedding() -> None:
    spec = {"$schema": "https://vega.github.io/schema/vega-lite/v5.json", "mark": "bar"}
    with patch.object(
        pptx_slides_service.vlc,
        "vegalite_to_png",
        return_value=MINI_PNG,
    ) as mock_png:
        path = generate_pptx_slides(
            "Chart",
            [{"layout": "content_basic", "title": "With chart", "_widget_id": "w1"}],
            widget_charts={"w1": spec},
        )
        mock_png.assert_called()
    prs = Presentation(str(path))
    # Picture shape present
    slide = prs.slides[0]
    pic_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pic_shapes) >= 1


def test_font_fallback() -> None:
    with patch.object(
        pptx_slides_service,
        "_font_available",
        side_effect=lambda name: name == "Arial",
    ):
        assert pptx_slides_service._resolve_font("FantasyFontXYZ") == "Arial"


def test_table_rendering() -> None:
    path = generate_pptx_slides(
        "Tbl",
        [
            {
                "layout": "title_only",
                "title": "Metrics",
                "table": {
                    "data": [
                        ["KPI", "Q1"],
                        ["A", "1"],
                        ["B", "2"],
                    ],
                },
            },
        ],
    )
    prs = Presentation(str(path))
    slide = prs.slides[0]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    tbl = tables[0].table
    assert len(tbl.rows) == 3
    assert len(tbl.columns) == 2
    assert tbl.cell(0, 0).text == "KPI"
    assert tbl.cell(2, 1).text == "2"


def test_widget_ref_keys_matches_expected() -> None:
    assert WIDGET_REF_KEYS == {"_widget_id", "_left_widget_id", "_right_widget_id"}


def test_fill_corporate_named_layout_closing_heading_body_fallback() -> None:
    tpl = (
        Path(__file__).resolve().parent.parent
        / "assets"
        / "databricks-corp-template.pptx"
    )
    if not tpl.is_file():
        pytest.skip("databricks-corp-template.pptx not in repo assets/")
    prs = Presentation(str(tpl))
    named = _resolve_named_layout(prs, "closing")
    assert named is not None
    slide = prs.slides.add_slide(named)
    spec = {
        "layout": "closing",
        "heading": "Questions",
        "body": "Thanks for joining.",
    }
    _fill_corporate_named_layout(slide, "closing", spec)
    combined = "".join(
        (s.text_frame.text or "")
        for s in slide.shapes
        if getattr(s, "has_text_frame", False)
    )
    assert "Questions" in combined
    assert "Thanks for joining." in combined


def test_build_slide_closing_subtitle_falls_back_to_body_or_text() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(_blank_slide_layout(prs))
    brand = _merge_brand(None)
    _build_slide(
        slide,
        prs,
        {"layout": "closing", "title": "Questions", "body": "Thanks"},
        brand,
        "Arial",
        {},
    )
    blob = _all_text_blob(prs)
    assert "Questions" in blob
    assert "Thanks" in blob


def test_generate_pptx_modern_two_column_with_template_preserves_columns_on_blank() -> (
    None
):
    tpl = (
        Path(__file__).resolve().parent.parent
        / "assets"
        / "databricks-corp-template.pptx"
    )
    if not tpl.is_file():
        pytest.skip("databricks-corp-template.pptx not in repo assets/")
    path = generate_pptx_slides(
        "Twocol",
        [
            {
                "layout": "two-column",
                "title": "Split",
                "left_header": "L",
                "left": ["a", "b"],
                "right_header": "R",
                "right": ["c", "d"],
            },
        ],
        pptx_template_path=tpl,
    )
    prs = Presentation(str(path))
    assert len(prs.slides) == 1
    slide_layout_name = prs.slides[0].slide_layout.name or ""
    assert "blank" not in slide_layout_name.lower()
    assert "2 column" in slide_layout_name.lower() or "two" in slide_layout_name.lower()

    blob = _all_text_blob(prs)
    for fragment in ("Split", "L", "R", "a", "b", "c", "d"):
        assert fragment in blob


def test_generate_pptx_title_still_uses_corporate_template() -> None:
    tpl = (
        Path(__file__).resolve().parent.parent
        / "assets"
        / "databricks-corp-template.pptx"
    )
    if not tpl.is_file():
        pytest.skip("databricks-corp-template.pptx not in repo assets/")
    path = generate_pptx_slides(
        "T",
        [
            {
                "layout": "title",
                "title": "Corporate title",
                "subtitle": "Corporate subtitle",
            },
        ],
        pptx_template_path=tpl,
    )
    prs = Presentation(str(path))
    layout_name_lc = (prs.slides[0].slide_layout.name or "").strip().lower()
    assert any(
        layout_name_lc == pattern.strip().lower()
        for pattern in _NAMED_LAYOUT_PATTERNS["title"]
    )


def test_generate_pptx_big_number_preserves_subtitle_and_text() -> None:
    path = generate_pptx_slides(
        "BN",
        [
            {
                "layout": "big-number",
                "number": "47%",
                "text": "Quarterly growth",
                "subtitle": "Up from 22% last quarter",
            },
        ],
    )
    assert path.exists()
    prs = Presentation(str(path))
    blob = _all_text_blob(prs)
    assert "47%" in blob
    assert "Quarterly growth" in blob
    assert "Up from 22% last quarter" in blob


def test_build_slide_modern_two_column_without_template_draws_title() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(_blank_slide_layout(prs))
    brand = _merge_brand(None)
    _build_slide(
        slide,
        prs,
        {
            "layout": "two-column",
            "title": "Modern title",
            "left": ["a"],
            "right": ["b"],
        },
        brand,
        "Arial",
        {},
    )
    blob = _all_text_blob(prs)
    assert "Modern title" in blob
