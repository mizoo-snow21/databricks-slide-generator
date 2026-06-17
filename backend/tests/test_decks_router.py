from __future__ import annotations

import zipfile
from io import BytesIO
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from main import app
from models import (
    ChartAugmentation,
    ChartDesign,
    ChartHighlight,
    Deck,
    DesignTokens,
    Template,
    TemplateBrand,
    TemplateCreate,
    WidgetInfo,
)
from routers.decks import (
    _check_placeholder_leak,
    _prerender_widget_chart_data_uris,
    _tokens_dict_from_template,
    deck_job_store,
    deck_repo,
    get_deck_service,
)
from routers.templates import templates_service
from services.brand_styles import CATEGORY_PALETTE_BY_PRESET
from services.deck_service import DeckService
from services.genie_service import GenieAnswer

TINY_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06"
    b"\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01"
    b"\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _bar_chart_genie_answer() -> GenieAnswer:
    return GenieAnswer(
        question="Revenue by month?",
        sql="SELECT month, v FROM t",
        columns=["month", "v"],
        rows=[{"month": "Jan", "v": 10}],
        status="ok",
    )


def _wait_job(job_id: str, timeout: float = 15.0):
    return deck_job_store.wait(job_id, timeout=timeout)


def _patch_ask_many(
    answers: list[GenieAnswer] | None = None,
    warnings: list[str] | None = None,
):
    ok = answers if answers is not None else [_bar_chart_genie_answer()]
    warn = warnings if warnings is not None else []
    return patch(
        "routers.decks.genie_service.ask_many",
        new_callable=AsyncMock,
        return_value=(ok, warn),
    )


@pytest.fixture(autouse=True)
def clear_deck_repo() -> None:
    deck_repo.clear()
    deck_job_store.clear()
    yield
    deck_repo.clear()
    deck_job_store.clear()


def test_list_decks_empty() -> None:
    client = TestClient(app)
    r = client.get("/api/decks")
    assert r.status_code == 200
    assert r.json() == [] or isinstance(r.json(), list)


def test_get_deck_404() -> None:
    client = TestClient(app)
    r = client.get("/api/decks/nonexistent")
    assert r.status_code == 404


def test_edit_html_for_missing_deck_404() -> None:
    client = TestClient(app)
    r = client.get("/api/decks/nonexistent/edit-html")
    assert r.status_code == 404


def test_edit_html_includes_slide_and_inspector_signature() -> None:
    html = (
        "<html><head></head><body>"
        '<section class="slide" data-slide-id="s1" data-osd-id="el-root">'
        '<h1 data-osd-id="el-h">Title</h1>'
        "</section></body></html>"
    )
    deck_repo.insert_deck(
        Deck(
            id="deck-with-slide",
            user_id="demo-user",
            template_id="t",
            genie_space_id="d",
            google_slides_template_id="gs",
            html_doc=html,
            design_tokens={},
            theme_markdown="",
        )
    )
    client = TestClient(app)
    r = client.get("/api/decks/deck-with-slide/edit-html")
    assert r.status_code == 200
    body = r.text
    assert 'class="slide"' in body
    assert "osd:ready" in body


def test_audit_deck_endpoint_404() -> None:
    client = TestClient(app)
    r = client.post("/api/decks/nonexistent-deck/audit")
    assert r.status_code == 404


def test_import_pptx_creates_deck() -> None:
    template = templates_service.create(
        TemplateCreate(
            name="Import Test Template",
            description="",
            google_slides_template_id="",
        ),
        user_id="",
    )
    p = Presentation()
    layout = p.slide_layouts[1] if len(p.slide_layouts) > 1 else p.slide_layouts[0]
    for title in ("Cover", "Body"):
        slide = p.slides.add_slide(layout)
        slide.shapes.title.text = title
    buf = BytesIO()
    p.save(buf)
    pptx_bytes = buf.getvalue()

    client = TestClient(app)
    files = {
        "file": (
            "deck.pptx",
            pptx_bytes,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    }
    data = {"template_id": template.id}
    r = client.post("/api/decks/import-pptx", files=files, data=data)
    assert r.status_code == 201, r.text
    payload = r.json()
    assert "id" in payload
    html = payload.get("html_doc", "")
    assert 'class="slide"' in html
    assert "slide-01" in html


def test_import_pptx_rejects_non_pptx() -> None:
    template = templates_service.create(
        TemplateCreate(
            name="Another Import Template",
            description="",
            google_slides_template_id="",
        ),
        user_id="",
    )
    client = TestClient(app)
    files = {"file": ("note.txt", b"hello", "text/plain")}
    data = {"template_id": template.id}
    r = client.post("/api/decks/import-pptx", files=files, data=data)
    assert r.status_code == 400


def test_outline_upload_doc_accepts_md() -> None:
    client = TestClient(app)
    files = {"file": ("brief.md", b"# Title\n\nHello", "text/markdown")}
    r = client.post("/api/decks/outline/upload-doc", files=files)
    assert r.status_code == 200
    payload = r.json()
    assert payload["reference_doc_name"] == "brief.md"
    assert "# Title" in payload["reference_doc"]
    assert "Hello" in payload["reference_doc"]


def test_outline_upload_doc_rejects_pdf() -> None:
    client = TestClient(app)
    files = {"file": ("x.pdf", b"%PDF-1.4", "application/pdf")}
    r = client.post("/api/decks/outline/upload-doc", files=files)
    assert r.status_code == 400


def test_check_placeholder_leak_detects_lorem() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(
            "ppt/slides/slide1.xml",
            '<?xml version="1.0"?><root><a:t>Lorem ipsum dolor</a:t></root>',
        )
    leaks = _check_placeholder_leak(buf.getvalue())
    assert len(leaks) >= 1
    assert any("lorem ipsum" in entry.lower() for entry in leaks)


def test_check_placeholder_leak_clean_returns_empty() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(
            "ppt/slides/slide1.xml",
            '<?xml version="1.0"?><root><a:t>Real content</a:t></root>',
        )
    assert _check_placeholder_leak(buf.getvalue()) == []


def test_check_placeholder_leak_ignores_shape_name_attributes() -> None:
    """Corporate templates have <p:nvSpPr name='Placeholder 2'> — not user text."""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(
            "ppt/slides/slide1.xml",
            '<?xml version="1.0"?><root>'
            '<p:nvSpPr name="Placeholder 2"/>'
            "<a:t>Real content</a:t>"
            "</root>",
        )
    assert _check_placeholder_leak(buf.getvalue()) == []


def test_export_pptx_placeholder_leak_raises_500(tmp_path) -> None:
    template = templates_service.create(
        TemplateCreate(
            name="Export Leak Template",
            description="",
            google_slides_template_id="",
        ),
        user_id="",
    )
    html_doc = (
        '<html><body><section class="slide" data-slide-id="s1">'
        "<h1>X</h1></section></body></html>"
    )
    deck_repo.insert_deck(
        Deck(
            id="deck-export-leak",
            user_id="demo-user",
            template_id=template.id,
            genie_space_id="",
            google_slides_template_id="",
            html_doc=html_doc,
            design_tokens={},
            theme_markdown="",
        )
    )
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(
            "ppt/slides/slide1.xml",
            '<?xml version="1.0"?><root><a:t>Lorem ipsum dolor</a:t></root>',
        )
    leaky = buf.getvalue()
    out_file = tmp_path / "exported.pptx"
    out_file.write_bytes(leaky)

    client = TestClient(app)
    with (
        patch("routers.decks._html_to_spec_slides", return_value=[]),
        patch(
            "routers.decks.generate_pptx_slides",
            return_value=str(out_file),
        ),
    ):
        response = client.get("/api/decks/deck-export-leak/export/pptx")
    assert response.status_code == 500
    assert "placeholder" in response.json()["detail"].lower()


def test_create_deck_high_quality_audit_failure_returns_503(
    stub_obo_dependency,
) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Audit Fail Template",
            description="",
            google_slides_template_id="tpl",
        ),
        user_id="u",
    )
    fake_deck = Deck(
        id="audit-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="tpl",
        html_doc=(
            '<html><body><section class="slide" data-slide-id="s">'
            "<h1>H</h1></section></body></html>"
        ),
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    mock_svc.audit_and_fix_deck.side_effect = RuntimeError("judge unavailable")

    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch(
            "routers.decks._prerender_widget_chart_data_uris",
            return_value=({}, {}),
        ),
    ):
        response = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": True,
            },
        )
        assert response.status_code == 202
        job_id = response.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "error"
    assert job.status_code == 503
    assert "High-quality audit failed" in (job.error or "")


def test_create_deck_applies_brand_palette_for_preset(stub_obo_dependency) -> None:
    """Preset-driven CATEGORY_PALETTE_BY_PRESET is passed into Vega-Lite config."""
    import services.vegalite_service as vegalite_mod

    created = templates_service.create(
        TemplateCreate(
            name="Corp Dark Palette Template",
            description="",
            google_slides_template_id="gs",
            preset_id="databricks-corp-dark",
        ),
        user_id="u",
    )

    fake_deck = Deck(
        id="palette-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    captured: list[dict] = []

    _real_convert = vegalite_mod.convert_widget_to_vegalite

    def _spy_convert(spec: dict, rows: list, **kwargs: object) -> dict:
        out = _real_convert(spec, rows, **kwargs)
        assert out is not None
        captured.append(out)
        return out

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch.object(
            vegalite_mod,
            "convert_widget_to_vegalite",
            side_effect=_spy_convert,
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            return_value=TINY_PNG,
        ),
    ):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202, r.text
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    assert captured, "convert_widget_to_vegalite should have been called"
    assert (
        captured[0]["config"]["range"]["category"]
        == CATEGORY_PALETTE_BY_PRESET["databricks-corp-dark"]
    )


def test_create_deck_appends_genie_warnings_for_failed_questions(
    stub_obo_dependency,
) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Genie Warn Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )

    fake_deck = Deck(
        id="genie-warn-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with (
        _patch_ask_many(
            answers=[_bar_chart_genie_answer()],
            warnings=["Slow question: timeout waiting for Genie response"],
        ),
        patch(
            "routers.decks._prerender_widget_chart_data_uris",
            return_value=({"q0": "data:image/png;base64,abc"}, {}),
        ),
    ):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?", "Slow question"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202, r.text
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    warnings = mock_svc.generate_deck.call_args.kwargs["chart_warnings"]
    assert any("Slow question" in w for w in warnings)
    assert mock_svc.generate_deck.call_args.kwargs["genie_space_id"] == "space-1"
    assert mock_svc.generate_deck.call_args.kwargs["questions"] == [
        "Revenue by month?",
        "Slow question",
    ]


def test_create_deck_all_genie_questions_failed_returns_502(
    stub_obo_dependency,
) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Genie Fail Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    client = TestClient(app)
    with _patch_ask_many(answers=[], warnings=["Q1: failed", "Q2: failed"]):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Q1", "Q2"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "error"
    assert job.status_code == 502
    assert "All questions failed" in (job.error or "")


def test_create_deck_appends_null_filter_warning_when_rows_dropped(
    stub_obo_dependency,
) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Null Rows Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )

    null_rows_answer = GenieAnswer(
        question="Revenue by month?",
        sql="SELECT month, v FROM t",
        columns=["month", "v"],
        rows=[{"month": None, "v": 1}, {"month": "Feb", "v": 2}],
        status="ok",
    )

    fake_deck = Deck(
        id="null-rows-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with (
        _patch_ask_many(answers=[null_rows_answer]),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            return_value=TINY_PNG,
        ),
    ):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202, r.text
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    warnings = mock_svc.generate_deck.call_args.kwargs["chart_warnings"]
    joined = " ".join(warnings).lower()
    assert "chart skipped" in joined or "null" in joined


def test_create_deck_applies_augmentation_when_llm_returns_one(
    stub_obo_dependency,
) -> None:
    import routers.decks as decks_rt
    import services.vegalite_service as vegalite_mod

    created = templates_service.create(
        TemplateCreate(
            name="Aug Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )

    fake_deck = Deck(
        id="aug-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    captured: list[dict] = []

    def capture_png(vl: dict) -> bytes:
        captured.append(vl)
        return TINY_PNG

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch.object(
            decks_rt._deck_llm,
            "augment_chart_specs_for_deck",
            return_value=[
                ChartAugmentation(
                    widget_id="q0",
                    highlight=ChartHighlight(field="month", values=["Jan"]),
                )
            ],
        ),
        patch.object(
            vegalite_mod,
            "convert_widget_to_vegalite",
            wraps=vegalite_mod.convert_widget_to_vegalite,
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            side_effect=capture_png,
        ),
    ):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202, r.text
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    assert captured
    spec0 = captured[0]
    enc0 = (
        spec0.get("layer", [{}])[0].get("encoding")
        if "layer" in spec0
        else spec0.get("encoding")
    ) or {}
    color_enc = enc0.get("color")
    assert color_enc is not None
    assert "condition" in color_enc


def test_create_deck_no_augmentation_when_llm_fails(
    stub_obo_dependency,
) -> None:
    import routers.decks as decks_rt

    created = templates_service.create(
        TemplateCreate(
            name="Aug Fail Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )

    fake_deck = Deck(
        id="aug-fail-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    captured: list[dict] = []

    def capture_png(vl: dict) -> bytes:
        captured.append(vl)
        return TINY_PNG

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch.object(
            decks_rt._deck_llm,
            "augment_chart_specs_for_deck",
            side_effect=RuntimeError("boom"),
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            side_effect=capture_png,
        ),
    ):
        r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert r.status_code == 202, r.text
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    assert captured
    assert "color" not in captured[0].get("encoding", {})


def test_create_deck_returns_job_then_completes(stub_obo_dependency) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Happy Path Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    fake_deck = Deck(
        id="happy-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    mock_svc.get_deck.return_value = fake_deck
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch(
            "routers.decks._prerender_widget_chart_data_uris",
            return_value=({}, {}),
        ),
    ):
        response = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert response.status_code == 202
        body = response.json()
        assert body["status"] == "running"
        job_id = body["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "done"
    assert job.deck_id == "happy-deck"

    poll = client.get(f"/api/decks/jobs/{job_id}")
    assert poll.status_code == 200
    poll_body = poll.json()
    assert poll_body["status"] == "done"
    assert poll_body["deck_id"] == "happy-deck"

    deck_resp = client.get("/api/decks/happy-deck")
    assert deck_resp.status_code == 200
    assert deck_resp.json()["id"] == "happy-deck"


def test_get_deck_job_unknown_returns_404() -> None:
    client = TestClient(app)
    r = client.get("/api/decks/jobs/does-not-exist")
    assert r.status_code == 404
    assert r.json()["detail"] == "Job not found"


def test_post_deck_outline_uses_genie_answers(stub_obo_dependency) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Outline Genie Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_outline.return_value = [
        {"layout": "title", "title": "Cover", "summary": "Intro", "notes": ""},
    ]
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with _patch_ask_many():
        r = client.post(
            "/api/decks/outline",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
            },
        )
        assert r.status_code == 202, r.text
        assert r.json()["status"] == "running"
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None and job.status == "done"
    widgets = mock_svc.generate_outline.call_args.kwargs["widgets"]
    assert len(widgets) == 1
    assert widgets[0].widget_id == "q0"
    assert widgets[0].query_result_summary
    assert widgets[0].row_count == 1


def test_outline_job_happy_path_poll_returns_slides(stub_obo_dependency) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Outline Poll Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_outline.return_value = [
        {"layout": "title", "title": "Cover", "summary": "Intro", "notes": ""},
        {"layout": "content", "title": "Body", "summary": "Details", "notes": "n"},
    ]
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with _patch_ask_many():
        r = client.post(
            "/api/decks/outline",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
            },
        )
        assert r.status_code == 202
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None and job.status == "done"

    poll = client.get(f"/api/decks/outline/jobs/{job_id}")
    assert poll.status_code == 200
    body = poll.json()
    assert body["status"] == "done"
    assert len(body["slides"]) == 2
    assert body["slides"][0]["title"] == "Cover"
    assert body["slides"][1]["title"] == "Body"


def test_get_outline_job_unknown_returns_404() -> None:
    client = TestClient(app)
    r = client.get("/api/decks/outline/jobs/does-not-exist")
    assert r.status_code == 404
    assert r.json()["detail"] == "Job not found"


def test_outline_and_deck_job_kind_isolation(stub_obo_dependency) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Kind Isolation Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    fake_deck = Deck(
        id="kind-deck",
        user_id="demo-user",
        template_id=created.id,
        genie_space_id="space-1",
        google_slides_template_id="gs",
        html_doc="<html><body></body></html>",
        design_tokens={},
        theme_markdown="",
        status="draft",
    )
    mock_svc = MagicMock(spec=DeckService)
    mock_svc.generate_deck.return_value = fake_deck
    mock_svc.generate_outline.return_value = [
        {"layout": "title", "title": "Cover", "summary": "Intro", "notes": ""},
    ]
    app.dependency_overrides[get_deck_service] = lambda: mock_svc

    client = TestClient(app)
    with (
        _patch_ask_many(),
        patch(
            "routers.decks._prerender_widget_chart_data_uris",
            return_value=({}, {}),
        ),
    ):
        deck_r = client.post(
            "/api/decks",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
                "high_quality": False,
            },
        )
        assert deck_r.status_code == 202
        deck_job_id = deck_r.json()["job_id"]
        _wait_job(deck_job_id)

        outline_r = client.post(
            "/api/decks/outline",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Revenue by month?"],
            },
        )
        assert outline_r.status_code == 202
        outline_job_id = outline_r.json()["job_id"]
        _wait_job(outline_job_id)

    assert client.get(f"/api/decks/outline/jobs/{deck_job_id}").status_code == 404
    assert client.get(f"/api/decks/jobs/{outline_job_id}").status_code == 404


def test_outline_all_genie_questions_failed_returns_502(stub_obo_dependency) -> None:
    created = templates_service.create(
        TemplateCreate(
            name="Outline Genie Fail Template",
            description="",
            google_slides_template_id="gs",
        ),
        user_id="u",
    )
    client = TestClient(app)
    with _patch_ask_many(answers=[], warnings=["Q1: failed", "Q2: failed"]):
        r = client.post(
            "/api/decks/outline",
            json={
                "template_id": created.id,
                "genie_space_id": "space-1",
                "questions": ["Q1", "Q2"],
            },
        )
        assert r.status_code == 202
        job_id = r.json()["job_id"]
        job = _wait_job(job_id)
    assert job is not None
    assert job.status == "error"
    assert job.status_code == 502
    assert "All questions failed" in (job.error or "")

    poll = client.get(f"/api/decks/outline/jobs/{job_id}")
    assert poll.status_code == 200
    assert poll.json()["status"] == "error"
    assert poll.json()["status_code"] == 502


def _template_without_tokens(*, theme: str) -> Template:
    return Template(
        id="tpl-tokens",
        name="Token Derivation Template",
        google_slides_template_id="gs-tokens",
        theme=theme,
        brand=TemplateBrand(
            primary="#123456",
            secondary="#888888",
            accent="#0066CC",
            text_light="#ffffff",
            text_dark="#111111",
            font="Noto Sans JP",
        ),
        tokens=None,
    )


def test_tokens_derived_light_uses_primary_as_accent_and_light_bg() -> None:
    derived = _tokens_dict_from_template(_template_without_tokens(theme="light"))
    assert derived["palette"]["bg"] == "#ffffff"
    assert derived["palette"]["text"] == "#111111"
    assert derived["palette"]["accent"] == "#123456"
    assert derived["fonts"]["display"] == "Noto Sans JP"
    assert derived["fonts"]["body"] == "Noto Sans JP"


def test_tokens_derived_dark_flips_bg_and_text() -> None:
    derived = _tokens_dict_from_template(_template_without_tokens(theme="dark"))
    assert derived["palette"]["bg"] == "#111111"
    assert derived["palette"]["text"] == "#ffffff"
    assert derived["palette"]["accent"] == "#123456"


def test_present_tokens_returned_verbatim() -> None:
    explicit_tokens = DesignTokens(
        palette={
            "bg": "#0a0a0a",
            "text": "#f5f5f5",
            "accent": "#ff00aa",
            "muted": "#444444",
        },
        fonts={"display": "Georgia", "body": "Arial"},
        typeScale={"hero": 180, "title": 72, "body": 32, "caption": 20},
        spacing={"padding": 96, "gap": 32},
        radius=8,
    )
    template = Template(
        id="tpl-explicit",
        name="Explicit Tokens Template",
        google_slides_template_id="gs-explicit",
        theme="light",
        brand=TemplateBrand(
            primary="#123456",
            text_light="#ffffff",
            text_dark="#111111",
            font="Noto Sans JP",
        ),
        tokens=explicit_tokens,
    )
    assert _tokens_dict_from_template(template) == explicit_tokens.model_dump()


def _many_category_rows(n: int = 10) -> list[dict[str, float | str]]:
    return [{"category": f"C{i}", "value": float(i + 1)} for i in range(n)]


def test_prerender_widgets_with_data_includes_field_stats() -> None:
    import routers.decks as decks_rt

    rows = _many_category_rows(10)
    widgets = [
        WidgetInfo(
            widget_id="q0",
            title="Rank",
            viz_type="auto",
            columns=["category", "value"],
            row_count=len(rows),
        )
    ]
    captured: list[dict] = []

    def capture_llm(widgets_with_data: list[dict], **kwargs: object) -> list:
        captured.extend(widgets_with_data)
        return []

    with (
        patch.object(
            decks_rt._deck_llm,
            "augment_chart_specs_for_deck",
            side_effect=capture_llm,
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            return_value=TINY_PNG,
        ),
    ):
        _prerender_widget_chart_data_uris(widgets, {"q0": rows})

    assert len(captured) == 1
    stats = captured[0].get("field_stats")
    assert isinstance(stats, dict)
    assert stats["category"]["distinct_count"] == 10
    assert stats["value"]["min"] == 1.0
    assert stats["value"]["max"] == 10.0


def test_prerender_applies_design_top_n_to_chart_data() -> None:
    import routers.decks as decks_rt

    rows = _many_category_rows(10)
    widgets = [
        WidgetInfo(
            widget_id="q0",
            title="Rank",
            viz_type="auto",
            columns=["category", "value"],
            row_count=len(rows),
        )
    ]
    render_calls: list[dict] = []

    def capture_render(vl: dict) -> bytes:
        render_calls.append(vl)
        return TINY_PNG

    with (
        patch.object(
            decks_rt._deck_llm,
            "augment_chart_specs_for_deck",
            return_value=[
                ChartAugmentation(
                    widget_id="q0",
                    design=ChartDesign(
                        chart_type="bar",
                        category_field="category",
                        value_field="value",
                        aggregate="sum",
                        sort="value_desc",
                        top_n=3,
                    ),
                )
            ],
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            side_effect=capture_render,
        ),
    ):
        charts, _errors = _prerender_widget_chart_data_uris(widgets, {"q0": rows})

    assert "q0" in charts
    assert render_calls
    cats = {r["category"] for r in render_calls[0]["data"]["values"]}
    assert len(cats) <= 3


def test_prerender_passes_rendered_rows_to_augmentation() -> None:
    import routers.decks as decks_rt
    import services.vegalite_service as vegalite_mod

    rows = _many_category_rows(10)
    widgets = [
        WidgetInfo(
            widget_id="q0",
            title="Rank",
            viz_type="auto",
            columns=["category", "value"],
            row_count=len(rows),
        )
    ]
    rows_passed_to_aug: list[list[dict]] = []
    orig_apply = vegalite_mod.apply_augmentation_to_spec

    def capture_apply(
        vl: dict,
        aug_rows: list[dict],
        aug: ChartAugmentation,
        **kwargs: object,
    ) -> dict:
        rows_passed_to_aug.append(list(aug_rows))
        return orig_apply(vl, aug_rows, aug, **kwargs)

    with (
        patch.object(
            decks_rt._deck_llm,
            "augment_chart_specs_for_deck",
            return_value=[
                ChartAugmentation(
                    widget_id="q0",
                    design=ChartDesign(
                        chart_type="bar",
                        category_field="category",
                        value_field="value",
                        aggregate="sum",
                        sort="value_desc",
                        top_n=3,
                    ),
                    highlight=ChartHighlight(field="category", values=["C9"]),
                )
            ],
        ),
        patch.object(
            vegalite_mod,
            "apply_augmentation_to_spec",
            side_effect=capture_apply,
        ),
        patch(
            "services.pptx_slides_service._render_chart_to_png",
            return_value=TINY_PNG,
        ),
    ):
        charts, _errors = _prerender_widget_chart_data_uris(widgets, {"q0": rows})

    assert "q0" in charts
    assert len(rows_passed_to_aug) == 1
    rendered_cats = {r["category"] for r in rows_passed_to_aug[0]}
    assert len(rendered_cats) <= 3
    assert "C9" in rendered_cats
