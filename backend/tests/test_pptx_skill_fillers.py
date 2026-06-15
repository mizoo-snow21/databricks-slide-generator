"""Smoke tests for corporate template slide fillers (per-layout PPTX export)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from services.pptx_slides_service import _NAMED_LAYOUT_PATTERNS, generate_pptx_slides

_CORP_TEMPLATE = (
    Path(__file__).resolve().parent.parent / "assets" / "databricks-corp-template.pptx"
)

_HYBRID_LAYOUTS = frozenset(
    {
        "agenda",
        "timeline",
        "icon-grid",
        "stat-row",
        "pros-cons",
        "comparison",
        "checklist",
        "logos",
    }
)


def _all_shape_text(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                tx = (sh.text_frame.text or "").strip()
                if tx:
                    parts.append(tx)
    return "\n".join(parts)


def _layout_matches(key: str, layout_name: str) -> bool:
    n = (layout_name or "").strip()
    if key in _HYBRID_LAYOUTS:
        return n.upper() == "CUSTOM"
    for pat in _NAMED_LAYOUT_PATTERNS.get(key, []):
        if not pat:
            continue
        if pat.lower() == n.lower() or pat in n:
            return True
    return False


_LAYOUT_SPECS: list[tuple[str, dict, list[str]]] = [
    (
        "title",
        {
            "layout": "title",
            "title": "PPTX Title",
            "subtitle": "Sub",
            "author": "A",
            "date": "2026",
        },
        ["PPTX Title", "Sub"],
    ),
    (
        "section",
        {"layout": "section", "title": "Sec A", "body": "Section eyebrow"},
        ["Sec A", "Section eyebrow"],
    ),
    (
        "content",
        {"layout": "content", "title": "Agenda topic", "bullets": ["x", "y"]},
        ["Agenda topic", "x", "y"],
    ),
    (
        "two-column",
        {
            "layout": "two-column",
            "title": "TC",
            "left_header": "L",
            "right_header": "R",
            "left": ["a"],
            "right": ["b"],
        },
        ["TC", "L", "R", "a", "b"],
    ),
    (
        "three-column",
        {
            "layout": "three-column",
            "title": "TH",
            "columns": [
                {"header": "A", "items": ["1"]},
                {"header": "B", "items": ["2"]},
                {"header": "C", "items": ["3"]},
            ],
        },
        ["TH", "A", "B", "C", "1", "2", "3"],
    ),
    (
        "big-number",
        {
            "layout": "big-number",
            "number": "99%",
            "text": "Growth",
            "subtitle": "QoQ",
        },
        ["99%", "Growth", "QoQ"],
    ),
    (
        "callout",
        {"layout": "callout", "text": "Bold words", "source": "Src"},
        ["Bold words", "Src"],
    ),
    (
        "quote",
        {"layout": "quote", "quote": "Q text", "attribution": "Person"},
        ["Q text", "Person"],
    ),
    (
        "closing",
        {"layout": "closing", "title": "Thanks"},
        ["Thanks"],
    ),
    (
        "two-column-icons",
        {
            "layout": "two-column-icons",
            "title": "I2",
            "columns": [
                {"header": "h1", "items": ["i1"]},
                {"header": "h2", "items": ["i2"]},
            ],
        },
        ["I2", "h1", "h2", "i1", "i2"],
    ),
    (
        "three-column-icons",
        {
            "layout": "three-column-icons",
            "title": "I3",
            "columns": [
                {"header": "a", "items": ["1"]},
                {"header": "b", "items": ["2"]},
                {"header": "c", "items": ["3"]},
            ],
        },
        ["I3", "a", "b", "c", "1", "2", "3"],
    ),
    (
        "cards",
        {
            "layout": "cards",
            "title": "Cd",
            "cards": [
                {"header": "c1", "items": ["u"]},
                {"header": "c2", "items": ["v"]},
                {"header": "c3", "items": ["w"]},
            ],
        },
        ["Cd", "c1", "c2", "c3", "u", "v", "w"],
    ),
    (
        "card-right",
        {
            "layout": "card-right",
            "title": "CR",
            "bullets": ["p"],
            "card_content": "diagram",
        },
        ["CR", "p", "diagram"],
    ),
    (
        "card-left",
        {
            "layout": "card-left",
            "title": "CL",
            "bullets": ["q"],
            "card_content": "pic",
        },
        ["CL", "q", "pic"],
    ),
    (
        "card-full",
        {"layout": "card-full", "title": "CF", "content": "wide body"},
        ["CF", "wide body"],
    ),
    (
        "one-column",
        {"layout": "one-column", "title": "OC", "content": "prose"},
        ["OC", "prose"],
    ),
    (
        "section-description",
        {
            "layout": "section-description",
            "title": "SD",
            "subtitle": "eyebrow",
            "description": "More text",
        },
        ["SD", "eyebrow", "More text"],
    ),
    (
        "agenda",
        {"layout": "agenda", "title": "Ag", "items": ["First", "Second"]},
        ["Ag", "First", "Second"],
    ),
    (
        "timeline",
        {
            "layout": "timeline",
            "title": "Tl",
            "steps": [
                {"title": "S1", "description": "D1"},
            ],
        },
        ["Tl", "S1", "D1"],
    ),
    (
        "icon-grid",
        {
            "layout": "icon-grid",
            "title": "Grid",
            "items": [{"icon": "★", "title": "F", "description": "desc"}],
        },
        ["Grid", "F", "desc"],
    ),
    (
        "stat-row",
        {
            "layout": "stat-row",
            "title": "Stats",
            "stats": [{"value": "3", "label": "KPI"}],
        },
        ["Stats", "3", "KPI"],
    ),
    (
        "pros-cons",
        {
            "layout": "pros-cons",
            "title": "PC",
            "pros": ["good"],
            "cons": ["bad"],
        },
        ["PC", "good", "bad"],
    ),
    (
        "comparison",
        {
            "layout": "comparison",
            "title": "Vs",
            "left_label": "A",
            "right_label": "B",
        },
        ["Vs", "vs.", "A", "B"],
    ),
    (
        "checklist",
        {
            "layout": "checklist",
            "title": "Chk",
            "items": [{"text": "Todo", "checked": True}],
        },
        ["Chk", "Todo"],
    ),
    (
        "logos",
        {
            "layout": "logos",
            "title": "Partners",
            "subtitle": "Trusted by",
            "logos": ["Acme", "Beta"],
        },
        ["Partners", "Trusted by", "Acme", "Beta"],
    ),
]


@pytest.mark.parametrize("layout_key,spec,fragments", _LAYOUT_SPECS)
def test_corporate_layout_smoke(
    layout_key: str, spec: dict, fragments: list[str]
) -> None:
    if not _CORP_TEMPLATE.is_file():
        pytest.skip("databricks-corp-template.pptx not in repo assets/")
    path = generate_pptx_slides("L", [spec], pptx_template_path=_CORP_TEMPLATE)
    prs = Presentation(str(path))
    assert len(prs.slides) == 1
    layout_name = prs.slides[0].slide_layout.name or ""
    assert _layout_matches(layout_key, layout_name), (
        f"expected layout contract for {layout_key}, got {layout_name!r}"
    )
    blob = _all_shape_text(prs)
    for frag in fragments:
        assert frag in blob, f"missing {frag!r} in {layout_key}; blob={blob!r}"
