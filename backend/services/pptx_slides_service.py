"""PPTX slide deck generation mirroring html_slides_service layout contracts."""

from __future__ import annotations

import json
import logging
import os
import platform
import sys
import re
import subprocess
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any

import vl_convert as vlc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

_FONT_DIR = Path(__file__).parent.parent / "assets" / "fonts"
if _FONT_DIR.is_dir():
    try:
        vlc.register_font_directory(str(_FONT_DIR))
        print(
            f"[fonts] registered {_FONT_DIR}",
            file=sys.stderr,
            flush=True,
        )
        bundled = sorted(p.name for p in _FONT_DIR.glob("*.ttf"))
        print(
            f"[fonts] bundled font files: {bundled}",
            file=sys.stderr,
            flush=True,
        )
    except Exception as exc:  # noqa: BLE001
        print(
            f"[fonts] register_font_directory failed: {type(exc).__name__}: {exc}",
            file=sys.stderr,
            flush=True,
        )

from services.html_slides_service import _columns_2col, _merge_brand, _normalize_layout

logger = logging.getLogger(__name__)

WIDGET_REF_KEYS = {"_widget_id", "_left_widget_id", "_right_widget_id"}

_FONT_FALLBACK = [
    "Noto Sans JP",
    "Noto Sans CJK JP",
    "Meiryo",
    "Yu Gothic",
    "Arial",
]

_WIDGET_ID_ORDER = ("_widget_id", "_left_widget_id", "_right_widget_id")

_LIGHT_BG = RGBColor(248, 250, 252)
_ALT_ROW_A = RGBColor(255, 255, 255)
_ALT_ROW_B = RGBColor(243, 244, 246)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    s = hex_str.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _darken_rgb(hex_str: str, factor: float = 0.5) -> RGBColor:
    s = hex_str.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    r, g, b = int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    return RGBColor(int(r * factor), int(g * factor), int(b * factor))


_fc_list_blob_cache: str | None = None


def _fc_list_font_blob() -> str:
    global _fc_list_blob_cache
    if _fc_list_blob_cache is not None:
        return _fc_list_blob_cache
    from shutil import which

    fc = which("fc-list")
    if not fc:
        _fc_list_blob_cache = ""
        return ""
    try:
        proc = subprocess.run(
            [fc],
            capture_output=True,
            text=True,
            timeout=20,
            check=False,
        )
        _fc_list_blob_cache = proc.stdout if proc.returncode == 0 else ""
    except (OSError, subprocess.TimeoutExpired):
        _fc_list_blob_cache = ""
    return _fc_list_blob_cache


def _font_available(font_name: str) -> bool:
    name = font_name.strip()
    if not name:
        return False
    blob = _fc_list_font_blob()
    if blob and name.lower() in blob.lower():
        return True
    key = name.lower().replace(" ", "")
    for root in _mac_font_roots():
        if not root.is_dir():
            continue
        try:
            for _dirpath, _dirnames, filenames in os.walk(root):
                for fn in filenames:
                    low = fn.lower()
                    if not low.endswith((".ttf", ".otf", ".ttc", ".otc", ".dfont")):
                        continue
                    stem = Path(fn).stem.lower().replace(" ", "")
                    if key in stem or stem in key:
                        return True
        except OSError:
            continue
    return False


def _mac_font_roots() -> list[Path]:
    if platform.system() != "Darwin":
        return []
    home = Path.home()
    return [
        Path("/Library/Fonts"),
        home / "Library/Fonts",
        Path("/System/Library/Fonts"),
        Path("/System/Library/Assets/com_apple_MobileAsset_Font3"),
    ]


def _resolve_font(brand_font: str) -> str:
    candidates = [brand_font] + [f for f in _FONT_FALLBACK if f != brand_font]
    for font_name in candidates:
        if _font_available(font_name):
            return font_name
    return "Arial"


_RESOLVED_IMPORT_FONT = _resolve_font(_merge_brand(None)["font"])
logger.info(
    "pptx_slides_service resolved default brand font to %s",
    _RESOLVED_IMPORT_FONT,
)


def _bullet_parts(bullets: Any) -> list[str]:
    if not bullets:
        return []
    if isinstance(bullets, str):
        return [p.strip() for p in re.split(r"[\n;]", bullets) if p.strip()]
    if isinstance(bullets, list):
        return [str(b).strip() for b in bullets if str(b).strip()]
    return []


def _set_bullet(paragraph: Any) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for el in list(p_pr):
        tag = el.tag.split("}")[-1]
        if tag in ("buNone", "buChar", "buAutoNum", "buFont", "buSzPct", "buSzPts"):
            p_pr.remove(el)
    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", "\u2022")
    p_pr.append(bu_char)


def _collect_ordered_widget_ids(slides: list[dict[str, Any]]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for slide in slides:
        for key in _WIDGET_ID_ORDER:
            if key not in WIDGET_REF_KEYS:
                continue
            ref = slide.get(key)
            if not ref:
                continue
            sid = str(ref)
            if sid not in seen:
                seen.add(sid)
                ordered.append(sid)
    return ordered


def _render_chart_to_png(
    vl_spec: dict[str, Any],
    width: int = 720,
    height: int = 460,
    scale: float = 2.0,
) -> bytes:
    """Render a Vega-Lite spec to PNG bytes.

    Defaults match _chart_base (720x460, ~1.56:1). With scale=2 the actual
    PNG is 1440x920 — fits two-column / card-left / three-column slide
    slots without letterbox padding.
    """
    spec = dict(vl_spec)
    if spec.get("width") in (None, "container"):
        spec["width"] = width
    if spec.get("height") in (None, "container"):
        spec["height"] = height
    return vlc.vegalite_to_png(vl_spec=json.dumps(spec), scale=scale)


def _prerender_widget_charts(
    slides: list[dict[str, Any]],
    widget_charts: dict[str, dict[str, Any]] | None,
) -> dict[str, bytes]:
    pngs: dict[str, bytes] = {}
    if not widget_charts:
        return pngs
    for wid in _collect_ordered_widget_ids(slides):
        spec = widget_charts.get(wid)
        if spec:
            pngs[wid] = _render_chart_to_png(spec)
    return pngs


# Named-layout mapping for slides whose corporate template has distinctive
# chrome. Order in each list = preference; first match across the loaded
# template's masters wins.
# Direct port of TSHuss/Databricks-slide-skill/scripts/generate-pptx.py
# LAYOUT_MAPPINGS — covers all 17 direct-named layouts.
_NAMED_LAYOUT_PATTERNS: dict[str, list[str]] = {
    # Structural (dark by default)
    "title": ["1_3 Title Slide B - Dark", "3 Title Slide B - Light", "TITLE"],
    "section": ["Content E - Power Statement 3", "SECTION_HEADER"],
    "callout": ["Content E - Power Statement 2_1", "MAIN_POINT"],
    "quote": ["Content E - Power Statement 2_1", "MAIN_POINT"],
    "closing": ["Z - Closing Dark", "Z - Closing Light"],
    # Content (light by default)
    "content": ["7 Content A - Basic", "TITLE_AND_BODY"],
    "two-column": ["9 Content B - 2 Column", "TITLE_AND_TWO_COLUMNS"],
    "three-column": ["11 Content C - 3 Column"],
    "big-number": ["Content E - Power Statement 1", "BIG_NUMBER"],
    # Master 1/2 layouts
    "two-column-icons": ["10 Content B - 2 Column w/ Icon Spot"],
    "three-column-icons": ["12 Content C - 3 Column w/ Icon Spot"],
    "cards": ["13 Content C - 3 Column Cards"],
    "card-right": ["14 Content D - Card Right"],
    "card-left": ["15 Content D - Card Left"],
    "card-full": ["16 Content D - Card Large"],
    "one-column": ["7 Content A - Basic", "ONE_COLUMN_TEXT"],
    "section-description": [
        "Content E - Power Statement 2",
        "SECTION_TITLE_AND_DESCRIPTION",
    ],
    "agenda": ["CUSTOM"],
    "timeline": ["CUSTOM"],
    "icon-grid": ["CUSTOM"],
    "stat-row": ["CUSTOM"],
    "pros-cons": ["CUSTOM"],
    "comparison": ["CUSTOM"],
    "checklist": ["CUSTOM"],
    "logos": ["CUSTOM"],
}

# Map modern HTML/spec layout keys to legacy _build_slide branch names when
# no corporate named layout matched (placeholder-free template).
_LAYOUT_ALIAS: dict[str, str] = {
    "two-column": "content_2col",
    "two-column-icons": "content_2col",
    "three-column": "content_3col",
    "three-column-icons": "content_3col",
    "section-description": "section_break",
    "section": "section_break",
    "quote": "quote_dark",
    "callout": "content_basic_dark",
    "content": "content_basic",
    "one-column": "content_basic",
    "cards": "content_3col",
    "card-left": "content_2col",
    "card-right": "content_2col",
    "card-full": "title_only",
}


_MODERN_STRUCTURED_FALLBACK: frozenset[str] = frozenset(
    {
        "agenda",
        "stat-row",
        "timeline",
        "icon-grid",
        "checklist",
        "pros-cons",
        "comparison",
        "logos",
    },
)


def _normalize_spec_for_three_column(spec: dict[str, Any]) -> dict[str, Any]:
    """Map modern `columns: [{header, items}]` or card grids into headerN/bodyN."""
    if spec.get("header1"):
        return spec
    cards = spec.get("cards")
    if isinstance(cards, list) and cards:
        out = dict(spec)
        for i, card in enumerate(cards[:3], start=1):
            if isinstance(card, dict):
                hh = card.get("header") or ""
                items = card.get("items") or []
                ctxt = card.get("content") or ""
                body = ctxt or (
                    "\n".join(str(x) for x in items if str(x).strip())
                    if isinstance(items, list)
                    else str(items)
                )
                out[f"header{i}"] = str(hh)
                out[f"body{i}"] = body.strip()
        return out
    cols = spec.get("columns")
    if not (isinstance(cols, list) and cols and isinstance(cols[0], dict)):
        return spec
    out = dict(spec)
    for i, col in enumerate(cols[:3], start=1):
        if not isinstance(col, dict):
            continue
        hh = col.get("header") or ""
        items = col.get("items") or []
        body = (
            "\n".join(str(x) for x in items if str(x).strip())
            if isinstance(items, list)
            else str(items)
        )
        ctxt = col.get("content")
        if ctxt:
            body = (body + "\n" + str(ctxt)).strip() if body else str(ctxt)
        out[f"header{i}"] = str(hh)
        out[f"body{i}"] = body.strip()
    return out


def _structured_bullets_for_modern_layout(
    layout: str, spec: dict[str, Any]
) -> list[str]:
    """Flatten structured modern fields into bullet lines for content_basic fallback."""
    lines: list[str] = []
    if layout == "agenda":
        for it in spec.get("items") or []:
            s = str(it).strip()
            if s:
                lines.append(s)
    elif layout == "stat-row":
        for st in spec.get("stats") or []:
            if isinstance(st, dict):
                v = st.get("value")
                lb = st.get("label")
                pv = "" if v in (None, "") else str(v).strip()
                plb = "" if lb in (None, "") else str(lb).strip()
                if pv and plb:
                    lines.append(f"{pv} — {plb}")
                elif pv or plb:
                    lines.append(pv or plb)
            elif str(st).strip():
                lines.append(str(st).strip())
    elif layout == "timeline":
        for step in spec.get("steps") or []:
            if isinstance(step, dict):
                t = str(step.get("title") or "").strip()
                d = str(step.get("description") or "").strip()
                if t and d:
                    lines.append(f"{t}: {d}")
                elif t or d:
                    lines.append(t or d)
            elif str(step).strip():
                lines.append(str(step).strip())
    elif layout == "icon-grid":
        for it in spec.get("items") or []:
            if isinstance(it, dict):
                parts = [
                    str(it.get("icon") or "").strip(),
                    str(it.get("title") or "").strip(),
                    str(it.get("description") or "").strip(),
                ]
                line = " · ".join(p for p in parts if p)
                if line:
                    lines.append(line)
            elif str(it).strip():
                lines.append(str(it).strip())
    elif layout == "checklist":
        for it in spec.get("items") or []:
            if isinstance(it, dict):
                tx = str(it.get("text") or "").strip()
                if not tx:
                    continue
                checked = it.get("checked")
                mark = "☑ " if checked else "☐ "
                lines.append(mark + tx)
            elif str(it).strip():
                lines.append(str(it).strip())
    elif layout == "pros-cons":
        ph = str(spec.get("pros_header") or "Pros")
        ch = str(spec.get("cons_header") or "Cons")
        lines.append(ph)
        for p in spec.get("pros") or []:
            if str(p).strip():
                lines.append(f"  + {str(p).strip()}")
        lines.append(ch)
        for c in spec.get("cons") or []:
            if str(c).strip():
                lines.append(f"  - {str(c).strip()}")
    elif layout == "comparison":
        ll = str(spec.get("left_label") or "").strip()
        rl = str(spec.get("right_label") or "").strip()
        if ll:
            lines.append(ll)
        if rl:
            lines.append(rl)
    elif layout == "logos":
        for logo in spec.get("logos") or []:
            if str(logo).strip():
                lines.append(str(logo).strip())
    return lines


def _resolve_named_layout(prs: Presentation, layout_name: str) -> Any:
    """Find a corporate template's named layout for our spec layout name.

    Returns the layout object or None if no match across all masters.
    """
    patterns = _NAMED_LAYOUT_PATTERNS.get((layout_name or "").lower()) or []
    if not patterns:
        return None
    candidates: list[tuple[str, Any]] = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            candidates.append(((layout.name or "").lower(), layout))
    for pattern in patterns:
        p = pattern.lower()
        for name, layout in candidates:
            if name == p:
                return layout
    return None


def _fill_placeholder(slide: Any, idx: int, text: str) -> bool:
    """Fill the placeholder at the given index with text. Returns True if found."""
    if not text:
        return False
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text_frame.text = text
            return True
    return False


def _fill_placeholder_bullets(slide: Any, idx: int, items: list[str]) -> bool:
    """Fill the placeholder at the given index with bulleted items.

    Mirrors the upstream skill's `fill_bullets` pattern (one line per item,
    relying on the template's own bullet styling rather than overriding it).
    """
    if not items:
        return False
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            tf.text = items[0] if items else ""
            for line in items[1:]:
                p = tf.add_paragraph()
                p.text = line
            return True
    return False


def _placeholders_by_type_sorted(slide: Any, ph_type: int) -> list[Any]:
    """Return placeholders of a given type, sorted left-to-right then top-to-bottom.

    Used for column layouts where the template has multiple BODY placeholders
    laid out side-by-side. Matches upstream skill's `get_placeholders_by_type`.
    """
    result = []
    for shape in slide.placeholders:
        if shape.placeholder_format.type == ph_type:
            result.append(shape)
    result.sort(key=lambda s: (s.left or 0, s.top or 0))
    return result


def _fill_corporate_named_layout(
    slide: Any,
    layout_name: str,
    spec: dict[str, Any],
    font_name: str = "Arial",
) -> None:
    """Fill placeholders on a slide whose layout came from the corporate template.

    Generic enough to handle title/subtitle/body for the common patterns; per-
    layout intelligence (which placeholder index is which) follows the upstream
    LAYOUT_MAPPINGS conventions where possible.
    """
    title = spec.get("title") or spec.get("text") or spec.get("quote") or ""
    subtitle = (
        spec.get("subtitle") or spec.get("source") or spec.get("attribution") or ""
    )
    if layout_name == "closing":
        if not title:
            title = spec.get("heading") or spec.get("h1") or ""
        if not subtitle:
            subtitle = spec.get("body") or spec.get("text") or ""

    bullets = spec.get("bullets") or []
    columns = list(spec.get("columns") or [])
    if layout_name in ("two-column", "two-column-icons") and not columns:
        left = spec.get("left") or []
        right = spec.get("right") or []
        if left or right:
            columns = [
                {
                    "header": spec.get("left_header") or "",
                    "items": left
                    if isinstance(left, list)
                    else [str(left)]
                    if left
                    else [],
                },
                {
                    "header": spec.get("right_header") or "",
                    "items": right
                    if isinstance(right, list)
                    else [str(right)]
                    if right
                    else [],
                },
            ]
    cards = spec.get("cards") or []

    # title placeholder is almost always idx 0
    if layout_name == "closing":
        title_ok = True
        if title:
            title_ok = _fill_placeholder(slide, 0, title)
        if subtitle and not _fill_placeholder(slide, 1, subtitle):
            _fill_placeholder(slide, 2, subtitle)

        placeholders_empty = not list(slide.placeholders)
        need_closing_textboxes = placeholders_empty or (bool(title) and not title_ok)
        if need_closing_textboxes:
            margin = Inches(1.0)
            content_w = Inches(11.333)
            white = RGBColor(0xFF, 0xFF, 0xFF)
            if title:
                _textbox(
                    slide,
                    margin,
                    Inches(2.7),
                    content_w,
                    Inches(1.0),
                    str(title),
                    font_name=font_name,
                    size_pt=34,
                    bold=True,
                    color=white,
                    align=PP_ALIGN.CENTER,
                )
            if subtitle:
                _textbox(
                    slide,
                    margin,
                    Inches(3.75),
                    content_w,
                    Inches(0.85),
                    str(subtitle),
                    font_name=font_name,
                    size_pt=18,
                    color=white,
                    align=PP_ALIGN.CENTER,
                )
    else:
        if title:
            _fill_placeholder(slide, 0, title)
        # subtitle placeholder commonly idx 1 or 2 (varies by layout)
        if subtitle and not _fill_placeholder(slide, 1, subtitle):
            _fill_placeholder(slide, 2, subtitle)

    # Body bullets for content / one-column layouts
    if bullets:
        # Try idx 1 first (typical body), then idx 2
        if not _fill_placeholder_bullets(slide, 1, bullets):
            _fill_placeholder_bullets(slide, 2, bullets)

    # Two/three-column layouts: fill BODY placeholders sorted by position.
    # PP_PLACEHOLDER.BODY = 2
    if columns:
        try:
            from pptx.enum.text import PP_PLACEHOLDER  # type: ignore[import-not-found]

            body_phs = _placeholders_by_type_sorted(slide, PP_PLACEHOLDER.BODY)
        except Exception:
            body_phs = []
        for i, col in enumerate(columns):
            if i >= len(body_phs):
                break
            ph = body_phs[i]
            tf = ph.text_frame
            header = col.get("header") or ""
            items = col.get("items") or []
            # Header as first paragraph (bold), items below
            tf.text = header
            for line in items:
                p = tf.add_paragraph()
                p.text = line

    # Cards layout: similar to columns
    if cards and not columns:
        try:
            from pptx.enum.text import PP_PLACEHOLDER  # type: ignore[import-not-found]

            body_phs = _placeholders_by_type_sorted(slide, PP_PLACEHOLDER.BODY)
        except Exception:
            body_phs = []
        for i, card in enumerate(cards):
            if i >= len(body_phs):
                break
            ph = body_phs[i]
            tf = ph.text_frame
            tf.text = card.get("header") or ""
            content = card.get("content") or ""
            items = card.get("items") or []
            if content:
                p = tf.add_paragraph()
                p.text = content
            for line in items:
                p = tf.add_paragraph()
                p.text = line

    # Big-number layout: value is the focal text
    if layout_name == "big-number":
        value = spec.get("value") or spec.get("number") or ""
        text = spec.get("text") or spec.get("label") or ""
        if value:
            _fill_placeholder(slide, 0, value)
        if text:
            _fill_placeholder(slide, 1, text)


def _blank_slide_layout(prs: Presentation) -> Any:
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name:
            return layout
    return prs.slide_layouts[6]


def _add_full_slide_rect(slide: Any, prs: Presentation) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    shape.line.fill.background()
    return shape


def _apply_primary_gradient_rect(shape: Any, primary_hex: str) -> None:
    try:
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = 135.0
        stops = fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = _hex_to_rgb(primary_hex)
        stops[1].position = 1.0
        stops[1].color.rgb = _darken_rgb(primary_hex, 0.45)
    except (AttributeError, ValueError, TypeError):
        _apply_solid_fill(shape, _hex_to_rgb(primary_hex))


def _apply_solid_fill(shape: Any, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_run_font(
    run: Any,
    *,
    font_name: str,
    size_pt: int | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    color: RGBColor | None = None,
) -> None:
    run.font.name = font_name
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _textbox(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    font_name: str,
    size_pt: int,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
) -> Any:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = valign
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        _set_run_font(
            run,
            font_name=font_name,
            size_pt=size_pt,
            bold=bold,
            italic=italic,
            color=color,
        )
    return box


def _accent_bar(slide: Any, left: Any, top: Any, width: Any, accent_hex: str) -> None:
    h = Inches(0.06)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h)
    bar.line.fill.background()
    _apply_solid_fill(bar, _hex_to_rgb(accent_hex))


def _add_bullets(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    parts: list[str],
    *,
    font_name: str,
    color: RGBColor,
    size_pt: int = 14,
) -> None:
    if not parts:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.text = part
        for run in p.runs:
            _set_run_font(run, font_name=font_name, size_pt=size_pt, color=color)
        _set_bullet(p)


def _add_table(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    rows: list[list[Any]],
    *,
    font_name: str,
    accent_hex: str,
    header_text: RGBColor,
    dark: bool,
) -> None:
    if not rows:
        return
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    graphic_frame = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = graphic_frame.table
    accent = _hex_to_rgb(accent_hex)
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            val = row[ci] if ci < len(row) else ""
            cell.text = str(val)
            if ri == 0:
                _apply_solid_fill(cell, accent)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        _set_run_font(
                            run,
                            font_name=font_name,
                            size_pt=11,
                            bold=True,
                            color=header_text,
                        )
            else:
                bg = _ALT_ROW_A if ri % 2 == 1 else _ALT_ROW_B
                _apply_solid_fill(cell, bg)
                tc = _hex_to_rgb("#202124") if not dark else _hex_to_rgb("#FFFFFF")
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        _set_run_font(run, font_name=font_name, size_pt=11, color=tc)


def _table_rows_from_slide(table: Any) -> list[list[Any]] | None:
    if not isinstance(table, dict):
        return None
    data = table.get("data")
    if not data or not isinstance(data, list):
        return None
    rows: list[list[Any]] = []
    for row in data:
        if isinstance(row, (list, tuple)):
            rows.append(list(row))
    return rows or None


def _place_picture(
    slide: Any,
    png: bytes,
    left: Any,
    top: Any,
    width: Any,
) -> None:
    stream = BytesIO(png)
    slide.shapes.add_picture(stream, left, top, width=width)


def _widget_pngs_for_slide(
    slide: dict[str, Any],
    chart_pngs: dict[str, bytes],
) -> tuple[bytes | None, bytes | None, bytes | None]:
    wid = slide.get("_widget_id")
    left_w = slide.get("_left_widget_id")
    right_w = slide.get("_right_widget_id")
    mid = chart_pngs.get(str(wid)) if wid else None
    lp = chart_pngs.get(str(left_w)) if left_w else None
    rp = chart_pngs.get(str(right_w)) if right_w else None
    return mid, lp, rp


def _build_slide(
    slide: Any,
    prs: Presentation,
    spec: dict[str, Any],
    brand: dict[str, str],
    font_name: str,
    chart_pngs: dict[str, bytes],
) -> None:
    raw_layout = _normalize_layout(str(spec.get("layout") or ""))
    layout = _LAYOUT_ALIAS.get(raw_layout, raw_layout)

    work = dict(spec)
    if raw_layout == "callout":
        if not work.get("title") and work.get("text"):
            work["title"] = work["text"]
        if not work.get("subtitle") and work.get("source"):
            work["subtitle"] = work["source"]
    if raw_layout in ("section", "section-description") and not work.get("subtitle"):
        work["subtitle"] = work.get("description") or work.get("body") or ""

    orig_layout_for_struct = layout
    if layout in _MODERN_STRUCTURED_FALLBACK:
        extra = _structured_bullets_for_modern_layout(orig_layout_for_struct, work)
        merged = _bullet_parts(work.get("bullets")) + extra
        work["bullets"] = merged
        layout = "content_basic"
    if layout == "content_3col":
        work = _normalize_spec_for_three_column(work)

    spec = work

    margin = Inches(0.65)
    content_w = prs.slide_width - 2 * margin
    primary = brand["primary"]
    accent = brand["accent"]
    text_dark = _hex_to_rgb(brand["text_dark"])
    text_light = _hex_to_rgb(brand["text_light"])
    secondary = _hex_to_rgb(brand["secondary"])

    mid_png, left_png, right_png = _widget_pngs_for_slide(spec, chart_pngs)

    dark_layouts = frozenset(
        {
            "section_break",
            "closing",
            "quote_dark",
            "content_basic_dark",
        },
    )
    is_dark = layout in dark_layouts or layout == "title"

    def light_bg_rect() -> None:
        bg = _add_full_slide_rect(slide, prs)
        _apply_solid_fill(bg, _LIGHT_BG)

    if layout == "title":
        bg = _add_full_slide_rect(slide, prs)
        _apply_primary_gradient_rect(bg, primary)
        _textbox(
            slide,
            margin,
            Inches(2.6),
            content_w,
            Inches(1.2),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=36,
            bold=True,
            color=text_light,
            align=PP_ALIGN.CENTER,
        )
        sub = spec.get("subtitle") or spec.get("body")
        if sub:
            _textbox(
                slide,
                margin,
                Inches(3.85),
                content_w,
                Inches(0.9),
                str(sub),
                font_name=font_name,
                size_pt=20,
                bold=False,
                color=text_light,
                align=PP_ALIGN.CENTER,
            )
        if mid_png:
            _place_picture(
                slide,
                mid_png,
                margin + Inches(1.0),
                Inches(5.0),
                Inches(6.0),
            )
        if left_png or right_png:
            y = Inches(5.0)
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, y, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    y,
                    half_w,
                )
        return

    if layout in ("content_basic", "content_subtitle"):
        light_bg_rect()
        top = Inches(0.45)
        _textbox(
            slide,
            margin,
            top,
            content_w,
            Inches(0.55),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=24,
            bold=True,
            color=text_dark,
        )
        _accent_bar(slide, margin, top + Inches(0.58), Inches(4.0), accent)
        y = Inches(1.15)
        sub = spec.get("subtitle")
        if sub:
            _textbox(
                slide,
                margin,
                y,
                content_w,
                Inches(0.45),
                str(sub),
                font_name=font_name,
                size_pt=16,
                color=secondary,
            )
            y += Inches(0.55)
        body = spec.get("body")
        if body:
            _textbox(
                slide,
                margin,
                y,
                content_w,
                Inches(0.85),
                str(body),
                font_name=font_name,
                size_pt=15,
                color=text_dark,
            )
            y += Inches(0.95)
        parts = _bullet_parts(spec.get("bullets"))
        if parts:
            _add_bullets(
                slide,
                margin,
                y,
                content_w,
                Inches(2.2),
                parts,
                font_name=font_name,
                color=text_dark,
            )
            y += Inches(2.25)
        if mid_png:
            _place_picture(slide, mid_png, margin, y, Inches(6.2))
        if left_png or right_png:
            y2 = y
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, y2, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    y2,
                    half_w,
                )
        return

    if layout == "content_2col":
        light_bg_rect()
        top = Inches(0.45)
        _textbox(
            slide,
            margin,
            top,
            content_w,
            Inches(0.55),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=24,
            bold=True,
            color=text_dark,
        )
        _accent_bar(slide, margin, top + Inches(0.58), Inches(4.0), accent)
        h1, b1, h2, b2 = _columns_2col(spec)
        col_w = content_w * 0.45
        gap = content_w * 0.10
        y = Inches(1.2)
        accent_c = _hex_to_rgb(accent)
        _textbox(
            slide,
            margin,
            y,
            col_w,
            Inches(0.4),
            h1,
            font_name=font_name,
            size_pt=16,
            bold=True,
            color=accent_c,
        )
        _textbox(
            slide,
            margin,
            y + Inches(0.42),
            col_w,
            Inches(2.5),
            b1,
            font_name=font_name,
            size_pt=14,
            color=text_dark,
        )
        left2 = margin + col_w + gap
        _textbox(
            slide,
            left2,
            y,
            col_w,
            Inches(0.4),
            h2,
            font_name=font_name,
            size_pt=16,
            bold=True,
            color=accent_c,
        )
        _textbox(
            slide,
            left2,
            y + Inches(0.42),
            col_w,
            Inches(2.5),
            b2,
            font_name=font_name,
            size_pt=14,
            color=text_dark,
        )
        yw = Inches(4.5)
        if mid_png:
            _place_picture(slide, mid_png, margin, yw, Inches(6.2))
        if left_png or right_png:
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, yw, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    yw,
                    half_w,
                )
        return

    if layout == "content_3col":
        light_bg_rect()
        top = Inches(0.45)
        _textbox(
            slide,
            margin,
            top,
            content_w,
            Inches(0.55),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=24,
            bold=True,
            color=text_dark,
        )
        _accent_bar(slide, margin, top + Inches(0.58), Inches(4.0), accent)
        y = Inches(1.15)
        col_w = content_w * 0.30
        gap = content_w * 0.05
        accent_c = _hex_to_rgb(accent)
        x = margin
        for i in (1, 2, 3):
            hh = str(spec.get(f"header{i}") or "")
            bb = str(spec.get(f"body{i}") or "")
            _textbox(
                slide,
                x,
                y,
                col_w,
                Inches(0.4),
                hh,
                font_name=font_name,
                size_pt=15,
                bold=True,
                color=accent_c,
            )
            _textbox(
                slide,
                x,
                y + Inches(0.42),
                col_w,
                Inches(2.4),
                bb,
                font_name=font_name,
                size_pt=13,
                color=text_dark,
            )
            x += col_w + gap
        yw = Inches(4.45)
        if mid_png:
            _place_picture(slide, mid_png, margin, yw, Inches(6.2))
        if left_png or right_png:
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, yw, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    yw,
                    half_w,
                )
        return

    if layout == "title_only":
        light_bg_rect()
        top = Inches(0.45)
        _textbox(
            slide,
            margin,
            top,
            content_w,
            Inches(0.55),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=24,
            bold=True,
            color=text_dark,
        )
        _accent_bar(slide, margin, top + Inches(0.58), Inches(4.0), accent)
        y = Inches(1.15)
        tbl_rows = _table_rows_from_slide(spec.get("table"))
        if tbl_rows:
            _add_table(
                slide,
                margin,
                y,
                content_w,
                Inches(3.2),
                tbl_rows,
                font_name=font_name,
                accent_hex=accent,
                header_text=text_light,
                dark=False,
            )
            y += Inches(3.35)
        if mid_png:
            _place_picture(slide, mid_png, margin, y, Inches(6.5))
        if left_png or right_png:
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, y, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    y,
                    half_w,
                )
        return

    if layout == "section_break":
        bg = _add_full_slide_rect(slide, prs)
        _apply_primary_gradient_rect(bg, primary)
        _textbox(
            slide,
            margin,
            Inches(2.9),
            content_w,
            Inches(1.0),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=32,
            bold=True,
            color=text_light,
            align=PP_ALIGN.CENTER,
        )
        sub = spec.get("subtitle") or spec.get("body") or spec.get("description")
        if sub:
            _textbox(
                slide,
                margin,
                Inches(4.0),
                content_w,
                Inches(0.7),
                str(sub),
                font_name=font_name,
                size_pt=18,
                color=text_light,
                align=PP_ALIGN.CENTER,
            )
        if mid_png:
            _place_picture(
                slide,
                mid_png,
                margin + Inches(1.0),
                Inches(4.6),
                Inches(6.0),
            )
        return

    if layout == "closing":
        bg = _add_full_slide_rect(slide, prs)
        _apply_primary_gradient_rect(bg, primary)
        _textbox(
            slide,
            margin,
            Inches(2.7),
            content_w,
            Inches(1.0),
            str(spec.get("title") or "Thank you"),
            font_name=font_name,
            size_pt=34,
            bold=True,
            color=text_light,
            align=PP_ALIGN.CENTER,
        )
        sub = spec.get("subtitle") or spec.get("body") or spec.get("text")
        if sub:
            _textbox(
                slide,
                margin,
                Inches(3.75),
                content_w,
                Inches(0.85),
                str(sub),
                font_name=font_name,
                size_pt=18,
                color=text_light,
                align=PP_ALIGN.CENTER,
            )
        if mid_png:
            _place_picture(
                slide,
                mid_png,
                margin + Inches(1.0),
                Inches(4.8),
                Inches(6.0),
            )
        return

    if layout == "quote_dark":
        bg = _add_full_slide_rect(slide, prs)
        _apply_primary_gradient_rect(bg, primary)
        quote = str(spec.get("title") or "")
        _textbox(
            slide,
            margin,
            Inches(2.5),
            content_w,
            Inches(1.4),
            quote,
            font_name=font_name,
            size_pt=22,
            bold=False,
            italic=True,
            color=text_light,
            align=PP_ALIGN.CENTER,
        )
        cite = spec.get("subtitle")
        if cite:
            _textbox(
                slide,
                margin,
                Inches(4.0),
                content_w,
                Inches(0.6),
                str(cite),
                font_name=font_name,
                size_pt=14,
                color=text_light,
                align=PP_ALIGN.CENTER,
            )
        if mid_png:
            _place_picture(
                slide,
                mid_png,
                margin + Inches(1.0),
                Inches(4.9),
                Inches(6.0),
            )
        return

    if layout == "content_basic_dark":
        bg = _add_full_slide_rect(slide, prs)
        _apply_primary_gradient_rect(bg, primary)
        top = Inches(0.45)
        _textbox(
            slide,
            margin,
            top,
            content_w,
            Inches(0.55),
            str(spec.get("title") or ""),
            font_name=font_name,
            size_pt=24,
            bold=True,
            color=text_light,
        )
        _accent_bar(slide, margin, top + Inches(0.58), Inches(4.0), accent)
        y = Inches(1.15)
        sub = spec.get("subtitle")
        if sub:
            _textbox(
                slide,
                margin,
                y,
                content_w,
                Inches(0.45),
                str(sub),
                font_name=font_name,
                size_pt=16,
                color=text_light,
            )
            y += Inches(0.55)
        body = spec.get("body")
        if body:
            _textbox(
                slide,
                margin,
                y,
                content_w,
                Inches(0.85),
                str(body),
                font_name=font_name,
                size_pt=15,
                color=text_light,
            )
            y += Inches(0.95)
        parts = _bullet_parts(spec.get("bullets"))
        if parts:
            _add_bullets(
                slide,
                margin,
                y,
                content_w,
                Inches(2.0),
                parts,
                font_name=font_name,
                color=text_light,
            )
            y += Inches(2.05)
        tbl_rows = _table_rows_from_slide(spec.get("table"))
        if tbl_rows:
            _add_table(
                slide,
                margin,
                y,
                content_w,
                Inches(2.6),
                tbl_rows,
                font_name=font_name,
                accent_hex=accent,
                header_text=text_light,
                dark=True,
            )
            y += Inches(2.7)
        if mid_png:
            _place_picture(slide, mid_png, margin, y, Inches(6.2))
        if left_png or right_png:
            half_w = content_w / 2 - Inches(0.15)
            if left_png:
                _place_picture(slide, left_png, margin, y, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    y,
                    half_w,
                )
        return

    if layout == "blank":
        light_bg_rect()
        if mid_png:
            _place_picture(
                slide,
                mid_png,
                margin + Inches(0.5),
                Inches(1.2),
                content_w - Inches(1.0),
            )
        elif left_png or right_png:
            half_w = content_w / 2 - Inches(0.15)
            y = Inches(1.2)
            if left_png:
                _place_picture(slide, left_png, margin, y, half_w)
            if right_png:
                _place_picture(
                    slide,
                    right_png,
                    margin + half_w + Inches(0.3),
                    y,
                    half_w,
                )
        return

    # Fallback (same spirit as html_slides_service)
    light_bg_rect()
    y = Inches(0.45)
    t = spec.get("title")
    if t:
        _textbox(
            slide,
            margin,
            y,
            content_w,
            Inches(0.55),
            str(t),
            font_name=font_name,
            size_pt=22,
            bold=True,
            color=text_dark,
        )
        y += Inches(0.75)
    b = spec.get("body")
    if b:
        _textbox(
            slide,
            margin,
            y,
            content_w,
            Inches(0.8),
            str(b),
            font_name=font_name,
            size_pt=14,
            color=text_dark,
        )
        y += Inches(0.9)
    parts = _bullet_parts(spec.get("bullets"))
    if parts:
        _add_bullets(
            slide,
            margin,
            y,
            content_w,
            Inches(2.0),
            parts,
            font_name=font_name,
            color=text_dark,
        )
        y += Inches(2.05)
    tbl_rows = _table_rows_from_slide(spec.get("table"))
    if tbl_rows:
        _add_table(
            slide,
            margin,
            y,
            content_w,
            Inches(2.5),
            tbl_rows,
            font_name=font_name,
            accent_hex=accent,
            header_text=text_light,
            dark=is_dark,
        )
        y += Inches(2.55)
    if mid_png:
        _place_picture(slide, mid_png, margin, y, Inches(6.0))
    if left_png or right_png:
        half_w = content_w / 2 - Inches(0.15)
        if left_png:
            _place_picture(slide, left_png, margin, y, half_w)
        if right_png:
            _place_picture(
                slide,
                right_png,
                margin + half_w + Inches(0.3),
                y,
                half_w,
            )


def generate_pptx_slides(
    title: str,
    slides: list[dict[str, Any]],
    brand: dict[str, Any] | None = None,
    widget_charts: dict[str, dict[str, Any]] | None = None,
    pptx_template_path: Path | None = None,
) -> Path:
    """Build a .pptx file from create-from-spec dicts; return path under /tmp.

    If pptx_template_path is provided, the new presentation inherits that file's
    slide masters / theme / fonts / footer. Any existing demo slides in the
    template are stripped before our generated slides are added.
    """
    _ = title  # deck metadata; file name only for now (parity with html title in doc)
    b = _merge_brand(brand)
    font_name = _resolve_font(b["font"])
    chart_pngs = _prerender_widget_charts(slides, widget_charts)

    if pptx_template_path is not None and pptx_template_path.is_file():
        prs = Presentation(str(pptx_template_path))
        # Strip any pre-existing demo slides AND their package parts so python-pptx
        # can re-allocate slide1.xml etc. without colliding (the upstream Databricks
        # template ships with 22 demo slides).
        sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
        parts_to_drop: list[Any] = []
        for sld in list(prs.slides):
            parts_to_drop.append(sld.part)
        # Remove from sldIdLst (deregisters the slide from the deck)
        for sld_id in list(sldIdLst):
            sldIdLst.remove(sld_id)
        # Drop relationships from the presentation part to each slide
        pres_part = prs.part
        rid_to_drop: list[str] = []
        for rid, rel in pres_part.rels.items():
            if rel.reltype.endswith("/slide"):
                rid_to_drop.append(rid)
        for rid in rid_to_drop:
            pres_part.drop_rel(rid)
        # Drop the slide parts themselves from the package
        package = pres_part.package
        for part in parts_to_drop:
            try:
                del package._parts[part.partname]  # type: ignore[attr-defined]
            except (AttributeError, KeyError):
                pass
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    layout_blank = _blank_slide_layout(prs)

    if not slides:
        slides = [{"layout": "title", "title": title or "", "subtitle": ""}]

    use_template = pptx_template_path is not None and pptx_template_path.is_file()
    corp_filler = None
    if use_template:
        from services.pptx_skill_fillers import CorporateSlideFiller, THEME_DICT

        deck_theme = {**THEME_DICT}
        if isinstance(b, dict):
            if b.get("primary"):
                deck_theme["dark_bg"] = b["primary"]
            if b.get("accent"):
                deck_theme["accent"] = b["accent"]
            if b.get("text_dark"):
                deck_theme["text_dark"] = b["text_dark"]
            if b.get("text_light"):
                deck_theme["text_light"] = b["text_light"]
        corp_filler = CorporateSlideFiller(prs, theme=deck_theme, font_name=font_name)
    for spec in slides:
        layout_name = (spec.get("layout") or "content").lower()
        if use_template:
            assert corp_filler is not None
            corp_filler.add_for_layout(layout_name, spec, chart_pngs=chart_pngs)
        else:
            spec_for_build = dict(spec)
            if layout_name == "big-number":
                num = spec_for_build.get("number") or spec_for_build.get("value")
                txt = spec_for_build.get("text") or spec_for_build.get("label") or ""
                sub = str(spec_for_build.get("subtitle") or "")
                title_cur = str(spec_for_build.get("title") or "").strip()
                body_cur = str(spec_for_build.get("body") or "").strip()
                if num and not title_cur:
                    spec_for_build["title"] = str(num)
                merged_body = "\n".join(p for p in (str(txt).strip(), sub.strip()) if p)
                if merged_body and not body_cur:
                    spec_for_build["body"] = merged_body
            slide = prs.slides.add_slide(layout_blank)
            _build_slide(slide, prs, spec_for_build, b, font_name, chart_pngs)

    out = Path(f"/tmp/genie-slide-{uuid.uuid4()}.pptx")
    prs.save(str(out))
    return out
