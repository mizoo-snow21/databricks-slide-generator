"""Corporate template slide fillers ported from slide-deck generate-pptx.py.

Fills Databricks corporate PPTX placeholders per layout and draws hybrid
content on the CUSTOM layout while preserving master chrome.
"""

from __future__ import annotations

import re
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.pptx_slides_service import (
    WIDGET_REF_KEYS,
    _NAMED_LAYOUT_PATTERNS,
    _hex_to_rgb,
    _place_picture,
    _widget_pngs_for_slide,
)

_LEGACY_PPTX_LAYOUT_ALIASES: dict[str, str] = {
    "content_basic": "content",
    "content_subtitle": "content",
    "section_break": "section",
    "title_only": "one-column",
    "quote_dark": "quote",
    "content_basic_dark": "callout",
    "content_2col": "two-column",
    "content_3col": "three-column",
}


THEME_DICT: dict[str, str] = {
    "accent": "#FF3621",
    "dark_bg": "#1B3139",
    "light_bg": "#F5F3F0",
    "text_dark": "#1B3139",
    "text_light": "#FFFFFF",
    "text_secondary": "#6B7280",
    "green": "#10B981",
    "red": "#EF4444",
    "divider": "#E5E7EB",
}

_CHART_FRIENDLY_LAYOUTS: frozenset[str] = frozenset(
    {
        "title_only",
        "two-column",
        "card-full",
        "card-left",
        "card-right",
        "big-number",
        "stat-row",
        "callout",
    }
)

# Double-asterisk (LLM / HTML deck) and single-asterisk (legacy skill) accents.
_ACCENT_PATTERN = re.compile(r"\*\*([^*]+)\*\*|\*([^*]+)\*")


def parse_accent_text(text: str) -> list[tuple[str, bool]]:
    """Split *text* / **text** into (segment, is_accent) tuples."""
    if not text:
        return [("", False)]
    if "*" not in text:
        return [(text, False)]
    segments: list[tuple[str, bool]] = []
    pos = 0
    for m in _ACCENT_PATTERN.finditer(text):
        if m.start() > pos:
            segments.append((text[pos : m.start()], False))
        inner = m.group(1) if m.group(1) is not None else m.group(2)
        segments.append((inner or "", True))
        pos = m.end()
    if pos < len(text):
        segments.append((text[pos:], False))
    return segments if segments else [(text, False)]


class CorporateSlideFiller:
    """Adaptation of slide-deck DatabricksSlideGenerator for an existing Presentation."""

    _DISPATCH: dict[str, str] = {
        "title": "add_title_slide",
        "section": "add_section_slide",
        "content": "add_content_slide",
        "two-column": "add_two_column_slide",
        "three-column": "add_three_column_slide",
        "big-number": "add_big_number_slide",
        "callout": "add_callout_slide",
        "quote": "add_quote_slide",
        "closing": "add_closing_slide",
        "two-column-icons": "add_two_column_icons_slide",
        "three-column-icons": "add_three_column_icons_slide",
        "cards": "add_cards_slide",
        "card-right": "add_card_right_slide",
        "card-left": "add_card_left_slide",
        "card-full": "add_card_full_slide",
        "one-column": "add_one_column_slide",
        "section-description": "add_section_description_slide",
        "agenda": "add_agenda_slide",
        "timeline": "add_timeline_slide",
        "icon-grid": "add_icon_grid_slide",
        "stat-row": "add_stat_row_slide",
        "pros-cons": "add_pros_cons_slide",
        "comparison": "add_comparison_slide",
        "checklist": "add_checklist_slide",
        "logos": "add_logos_slide",
    }

    def __init__(
        self,
        prs: Presentation,
        *,
        theme: dict[str, str] | None = None,
        font_name: str = "Arial",
    ) -> None:
        self.prs = prs
        self.theme = dict(THEME_DICT)
        if theme:
            self.theme.update(theme)
        self.font_name = font_name
        self.slide_count = 0

        self.layouts: dict[str, Any] = {}
        self.light_layouts: dict[str, Any] = {}
        self.dark_layouts: dict[str, Any] = {}

        for master in self.prs.slide_masters:
            is_dark = self._is_dark_background(master)
            for layout in master.slide_layouts:
                name = layout.name or ""
                if is_dark:
                    self.dark_layouts[name] = layout
                else:
                    self.light_layouts[name] = layout
                if name not in self.layouts:
                    self.layouts[name] = layout

    def _is_dark_background(self, master: Any) -> bool:
        try:
            fill = master.background.fill
            if fill.type is not None:
                if hasattr(fill, "fore_color") and hasattr(fill.fore_color, "rgb"):
                    rgb = fill.fore_color.rgb
                    if rgb:
                        rgb_str = str(rgb).upper()
                        if rgb_str == "1B3139":
                            return True
                        try:
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            if (r + g + b) < 384:
                                return True
                        except (ValueError, IndexError):
                            pass
        except (AttributeError, TypeError):
            pass
        return False

    def get_layout(self, slide_type: str, *, prefer_dark: bool = False) -> Any:
        patterns = _NAMED_LAYOUT_PATTERNS.get(slide_type, ["BLANK"])
        if prefer_dark:
            search_caches = [self.dark_layouts, self.layouts]
        else:
            search_caches = [self.layouts]
        for pattern in patterns:
            p = pattern
            for cache in search_caches:
                if p in cache:
                    return cache[p]
            for cache in search_caches:
                for name, layout in cache.items():
                    if p in name:
                        return layout
        return self.layouts.get("BLANK", next(iter(self.layouts.values())))

    def get_placeholder(
        self, slide: Any, idx: int | None = None, ph_type: int | None = None
    ) -> Any:
        for shape in slide.placeholders:
            if idx is not None and shape.placeholder_format.idx == idx:
                return shape
            if ph_type is not None and shape.placeholder_format.type == ph_type:
                return shape
        return None

    def get_placeholders_by_type(self, slide: Any, ph_type: int) -> list[Any]:
        matching = []
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                matching.append(shape)
        matching.sort(key=lambda s: (s.top, s.left))
        return matching

    def fill_text(
        self,
        placeholder: Any,
        text: str,
        font_size: int | None = None,
        bold: bool | None = None,
        color: str | None = None,
        align: PP_ALIGN | None = None,
    ) -> None:
        if placeholder is None:
            return
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        segments = parse_accent_text(text or "")
        has_accent = any(is_accent for _, is_accent in segments)
        accent_hex = self.theme["accent"]

        if has_accent:
            for i, (segment_text, is_accent) in enumerate(segments):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                run.text = segment_text
                run.font.name = self.font_name
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if is_accent:
                    run.font.color.rgb = _hex_to_rgb(accent_hex)
                elif color:
                    run.font.color.rgb = _hex_to_rgb(color)
        else:
            p.text = text or ""
            p.font.name = self.font_name
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = _hex_to_rgb(color)

    def fill_bullets(
        self, placeholder: Any, items: list[str], font_size: int | None = None
    ) -> None:
        if placeholder is None or not items:
            return
        tf = placeholder.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.level = 0
            p.font.name = self.font_name
            if font_size:
                p.font.size = Pt(font_size)

    def add_textbox(
        self,
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 18,
        bold: bool = False,
        color: str | None = None,
        alignment: int | None = None,
    ) -> None:
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        segments = parse_accent_text(text or "")
        has_accent = any(a for _, a in segments)
        accent_hex = self.theme["accent"]
        if has_accent:
            for i, (seg, is_acc) in enumerate(segments):
                if i == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                else:
                    run = p.add_run()
                run.text = seg
                run.font.name = self.font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if is_acc:
                    run.font.color.rgb = _hex_to_rgb(accent_hex)
                elif color:
                    run.font.color.rgb = _hex_to_rgb(color)
        else:
            p.text = text or ""
            p.font.name = self.font_name
            p.font.size = Pt(font_size)
            p.font.bold = bold
            if color:
                p.font.color.rgb = _hex_to_rgb(color)
        if alignment is not None:
            p.alignment = alignment  # type: ignore[assignment]

    def _create_slide(
        self,
        slide_type: str,
        data: dict[str, Any],
        *,
        prefer_dark: bool = False,
    ) -> Any:
        self.slide_count += 1
        layout = self.get_layout(slide_type, prefer_dark=prefer_dark)
        slide = self.prs.slides.add_slide(layout)
        notes = data.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
        return slide

    # --- Normalization (matches html/spec → skill JSON shape) ---

    @staticmethod
    def _normalize_spec(layout_name: str, spec: dict[str, Any]) -> dict[str, Any]:
        data = dict(spec)
        ln = layout_name
        if ln == "callout":
            if not data.get("title") and data.get("text"):
                data["title"] = data["text"]
            if not data.get("subtitle") and data.get("source"):
                data["subtitle"] = data["source"]
        if ln in ("section", "section-description") and not data.get("subtitle"):
            data["subtitle"] = data.get("description") or data.get("body") or ""
        if ln == "closing":
            if not data.get("title"):
                data["title"] = data.get("heading") or data.get("h1") or ""
            if not data.get("subtitle"):
                data["subtitle"] = data.get("body") or data.get("text") or ""
        if ln in ("two-column", "two-column-icons") and not data.get("columns"):
            if any(
                data.get(k) is not None and data.get(k) != ""
                for k in ("col1_header", "col1_body", "col2_header", "col2_body")
            ):
                data = dict(data)

                def _body_lines(v: Any) -> list[str]:
                    if v is None:
                        return []
                    if isinstance(v, list):
                        return [str(i) for i in v if str(i).strip()]
                    s = str(v).strip()
                    if not s:
                        return []
                    lines = [p.strip() for p in s.splitlines() if p.strip()]
                    return lines or [s]

                data.setdefault(
                    "left_header",
                    str(data.get("left_header") or data.get("col1_header") or ""),
                )
                data.setdefault(
                    "right_header",
                    str(data.get("right_header") or data.get("col2_header") or ""),
                )
                if not data.get("left"):
                    data["left"] = _body_lines(data.get("col1_body"))
                if not data.get("right"):
                    data["right"] = _body_lines(data.get("col2_body"))
            left = data.get("left") or []
            right = data.get("right") or []
            if left or right or data.get("left_header") or data.get("right_header"):
                data["columns"] = [
                    {
                        "header": data.get("left_header") or "",
                        "items": left
                        if isinstance(left, list)
                        else [str(left)]
                        if left
                        else [],
                    },
                    {
                        "header": data.get("right_header") or "",
                        "items": right
                        if isinstance(right, list)
                        else [str(right)]
                        if right
                        else [],
                    },
                ]
        return data

    def _place_charts(
        self,
        slide: Any,
        layout_name: str,
        spec: dict[str, Any],
        chart_pngs: dict[str, bytes] | None,
    ) -> None:
        if not chart_pngs:
            return
        mid, left_p, right_p = _widget_pngs_for_slide(spec, chart_pngs)
        prs = self.prs
        margin = Inches(0.65)
        content_w = prs.slide_width - 2 * margin
        if layout_name == "title":
            if mid:
                _place_picture(slide, mid, margin + Inches(1.0), Inches(5.0), Inches(6.0))
            if left_p or right_p:
                y = Inches(5.0)
                half_w = content_w / 2 - Inches(0.15)
                if left_p:
                    _place_picture(slide, left_p, margin, y, half_w)
                if right_p:
                    _place_picture(
                        slide,
                        right_p,
                        margin + half_w + Inches(0.3),
                        y,
                        half_w,
                    )
            return
        if mid and not left_p and not right_p:
            _place_picture(slide, mid, margin + Inches(0.5), Inches(5.0), Inches(7.0))
        elif left_p or right_p:
            y = Inches(5.2)
            half_w = content_w / 2 - Inches(0.15)
            if left_p:
                _place_picture(slide, left_p, margin, y, half_w)
            if right_p:
                _place_picture(
                    slide,
                    right_p,
                    margin + half_w + Inches(0.3),
                    y,
                    half_w,
                )

    def add_for_layout(
        self,
        layout_name: str,
        spec: dict[str, Any],
        *,
        chart_pngs: dict[str, bytes] | None = None,
    ) -> None:
        raw_in = (layout_name or "content").lower()
        raw = _LEGACY_PPTX_LAYOUT_ALIASES.get(raw_in, raw_in)
        if raw not in _NAMED_LAYOUT_PATTERNS:
            raw = "content"
        data = self._normalize_spec(raw, spec)
        method = self._DISPATCH.get(raw, "add_content_slide")
        getattr(self, method)(data)
        last = self.prs.slides[-1]
        if raw in _CHART_FRIENDLY_LAYOUTS or any(
            k in data for k in WIDGET_REF_KEYS
        ):
            self._place_charts(last, raw, data, chart_pngs)

    # --- Ported add_* methods (skill bodies) ---

    def add_title_slide(self, data: dict[str, Any], *, prefer_dark: bool = True) -> None:
        slide = self._create_slide("title", data, prefer_dark=prefer_dark)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Presentation Title"))
        subtitle_ph = self.get_placeholder(slide, idx=1)
        subtitle_parts = []
        if data.get("subtitle"):
            subtitle_parts.append(str(data["subtitle"]))
        self.fill_text(
            subtitle_ph, "\n".join(subtitle_parts) if subtitle_parts else ""
        )
        author_ph = self.get_placeholder(slide, idx=2)
        if author_ph:
            author_parts = []
            if data.get("author"):
                author_parts.append(str(data["author"]))
            if data.get("date"):
                author_parts.append(str(data["date"]))
            self.fill_text(
                author_ph, " | ".join(author_parts) if author_parts else ""
            )

    def add_section_slide(self, data: dict[str, Any], *, prefer_dark: bool = True) -> None:
        slide = self._create_slide("section", data, prefer_dark=prefer_dark)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))
        sub = str(data.get("subtitle") or "").strip()
        if sub:
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, sub)
            else:
                self.add_textbox(
                    slide,
                    sub,
                    0.65,
                    4.5,
                    12.0,
                    1.5,
                    font_size=18,
                    color=self.theme["text_light"],
                )

    def add_content_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("content", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Slide Title"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        body_ph = self.get_placeholder(slide, idx=1)
        bullets = data.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [p.strip() for p in re.split(r"[\n;]", bullets) if p.strip()]
        self.fill_bullets(body_ph, [str(b) for b in bullets])

    def add_two_column_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("two-column", data)
        title_ph = self.get_placeholder(slide, ph_type=PP_PLACEHOLDER.TITLE)
        self.fill_text(title_ph, data.get("title", "Two Column"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        subtitle_phs = self.get_placeholders_by_type(slide, PP_PLACEHOLDER.SUBTITLE)
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)
        body_phs = self.get_placeholders_by_type(slide, PP_PLACEHOLDER.BODY)
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)
        columns = data.get("columns") or []
        if columns:
            for i, col in enumerate(columns[:2]):
                if not isinstance(col, dict):
                    continue
                if i < len(col_headers) and col.get("header"):
                    self.fill_text(col_headers[i], str(col["header"]))
                if i < len(col_bodies):
                    items = col.get("items") or []
                    self.fill_bullets(col_bodies[i], [str(x) for x in items])
        else:
            if data.get("left_header") and len(col_headers) > 0:
                self.fill_text(col_headers[0], str(data["left_header"]))
            if data.get("right_header") and len(col_headers) > 1:
                self.fill_text(col_headers[1], str(data["right_header"]))
            if len(col_bodies) > 0:
                left_items = data.get("left") or []
                if not isinstance(left_items, list):
                    left_items = [str(left_items)] if left_items else []
                self.fill_bullets(col_bodies[0], [str(x) for x in left_items])
            if len(col_bodies) > 1:
                right_items = data.get("right") or []
                if not isinstance(right_items, list):
                    right_items = [str(right_items)] if right_items else []
                self.fill_bullets(col_bodies[1], [str(x) for x in right_items])

    def add_three_column_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("three-column", data)
        title_ph = self.get_placeholder(slide, ph_type=PP_PLACEHOLDER.TITLE)
        self.fill_text(title_ph, data.get("title", "Three Column"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        columns = data.get("columns", [])
        subtitle_phs = self.get_placeholders_by_type(slide, PP_PLACEHOLDER.SUBTITLE)
        col_headers = [ph for ph in subtitle_phs if ph.top.inches > 1.5]
        col_headers.sort(key=lambda s: s.left)
        body_phs = self.get_placeholders_by_type(slide, PP_PLACEHOLDER.BODY)
        col_bodies = [ph for ph in body_phs if ph.top.inches > 2.5]
        col_bodies.sort(key=lambda s: s.left)
        for i, col in enumerate(columns[:3]):
            if not isinstance(col, dict):
                continue
            if i < len(col_headers) and col.get("header"):
                self.fill_text(col_headers[i], str(col["header"]))
            if i < len(col_bodies):
                items = col.get("items") or []
                self.fill_bullets(col_bodies[i], [str(x) for x in items])

    def add_big_number_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("big-number", data)
        title_ph = self.get_placeholder(slide, idx=0)
        num = str(data.get("number", data.get("value", "0")))
        self.fill_text(
            title_ph, num, bold=True, color=self.theme["accent"]
        )
        body_ph = self.get_placeholder(slide, idx=1)
        description = str(data.get("text", "") or data.get("label", "") or "")
        if data.get("subtitle"):
            description += f"\n{data['subtitle']}"
        if body_ph is not None:
            self.fill_text(body_ph, description)
        elif description:
            self.add_textbox(
                slide,
                description,
                left=1.0,
                top=5.5,
                width=11.33,
                height=1.2,
                font_size=20,
                color=self.theme["text_light"],
                alignment=PP_ALIGN.CENTER,
            )

    def add_callout_slide(self, data: dict[str, Any], *, prefer_dark: bool = True) -> None:
        slide = self._create_slide("callout", data, prefer_dark=prefer_dark)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, str(data.get("text", "Key message")))
        if data.get("source"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['source']}")

    def add_quote_slide(self, data: dict[str, Any], *, prefer_dark: bool = True) -> None:
        slide = self._create_slide("quote", data, prefer_dark=prefer_dark)
        quote_text = f'"{data.get("quote", "Quote goes here.")}"'
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, quote_text)
        if data.get("attribution"):
            subtitle_ph = self.get_placeholder(slide, idx=1)
            if subtitle_ph:
                self.fill_text(subtitle_ph, f"— {data['attribution']}")

    def add_closing_slide(self, data: dict[str, Any], *, prefer_dark: bool = True) -> None:
        slide = self._create_slide("closing", data, prefer_dark=prefer_dark)
        text_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        title = str(data.get("title", "Thank You") or "Thank You")
        self.add_textbox(
            slide,
            title,
            0.75,
            0.8,
            11.5,
            1.2,
            font_size=48,
            bold=True,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
        )
        sub = str(data.get("subtitle") or "").strip()
        if sub:
            self.add_textbox(
                slide,
                sub,
                0.75,
                2.0,
                11.5,
                0.85,
                font_size=18,
                bold=False,
                color=text_color,
                alignment=PP_ALIGN.CENTER,
            )

    def add_two_column_icons_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("two-column-icons", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Two Column with Icons"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=5)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        columns = data.get("columns", [])
        for i, col in enumerate(columns[:2]):
            if not isinstance(col, dict):
                continue
            header_ph = self.get_placeholder(slide, idx=3 + i)
            if header_ph and col.get("header"):
                self.fill_text(header_ph, str(col["header"]))
        for i, col in enumerate(columns[:2]):
            if not isinstance(col, dict):
                continue
            body_ph = self.get_placeholder(slide, idx=1 + i)
            items = col.get("items") or []
            self.fill_bullets(body_ph, [str(x) for x in items])

    def add_three_column_icons_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("three-column-icons", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Three Column with Icons"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        columns = data.get("columns", [])
        header_indices = [3, 4, 6]
        for i, col in enumerate(columns[:3]):
            if not isinstance(col, dict):
                continue
            header_ph = self.get_placeholder(slide, idx=header_indices[i])
            if header_ph and col.get("header"):
                self.fill_text(header_ph, str(col["header"]))
        body_indices = [1, 2, 5]
        for i, col in enumerate(columns[:3]):
            if not isinstance(col, dict):
                continue
            body_ph = self.get_placeholder(slide, idx=body_indices[i])
            items = col.get("items") or []
            self.fill_bullets(body_ph, [str(x) for x in items])

    def add_cards_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("cards", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Cards"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=7)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        cards = data.get("cards", [])
        for i, card in enumerate(cards[:3]):
            if not isinstance(card, dict):
                continue
            header_ph = self.get_placeholder(slide, idx=4 + i)
            if header_ph and card.get("header"):
                self.fill_text(header_ph, str(card["header"]))
        for i, card in enumerate(cards[:3]):
            if not isinstance(card, dict):
                continue
            body_ph = self.get_placeholder(slide, idx=1 + i)
            if card.get("content"):
                self.fill_text(body_ph, str(card["content"]))
            elif card.get("items"):
                self.fill_bullets(
                    body_ph, [str(x) for x in card.get("items") or []]
                )

    def add_card_right_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("card-right", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Right"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        left_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(left_ph, str(data["content"]))
        elif data.get("bullets"):
            b = data["bullets"]
            if isinstance(b, str):
                b = [p.strip() for p in re.split(r"[\n;]", b) if p.strip()]
            self.fill_bullets(left_ph, [str(x) for x in b])
        right_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(right_ph, str(data["card_content"]))

    def add_card_left_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("card-left", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Card Left"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=3)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        left_ph = self.get_placeholder(slide, idx=2)
        if data.get("card_content"):
            self.fill_text(left_ph, str(data["card_content"]))
        right_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(right_ph, str(data["content"]))
        elif data.get("bullets"):
            b = data["bullets"]
            if isinstance(b, str):
                b = [p.strip() for p in re.split(r"[\n;]", b) if p.strip()]
            self.fill_bullets(right_ph, [str(x) for x in b])

    def add_card_full_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("card-full", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Full Card"))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        card_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(card_ph, str(data["content"]))

    def add_one_column_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("one-column", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, str(data.get("title", "")))
        if data.get("subtitle"):
            subtitle_ph = self.get_placeholder(slide, idx=2)
            if subtitle_ph:
                self.fill_text(subtitle_ph, str(data["subtitle"]))
        body_ph = self.get_placeholder(slide, idx=1)
        if data.get("content"):
            self.fill_text(body_ph, str(data["content"]))
        elif data.get("bullets"):
            b = data["bullets"]
            if isinstance(b, str):
                b = [p.strip() for p in re.split(r"[\n;]", b) if p.strip()]
            self.fill_bullets(body_ph, [str(x) for x in b])

    def add_section_description_slide(self, data: dict[str, Any]) -> None:
        slide = self._create_slide("section-description", data)
        title_ph = self.get_placeholder(slide, idx=0)
        self.fill_text(title_ph, data.get("title", "Section Title"))
        subtitle_ph = self.get_placeholder(slide, idx=1)
        if subtitle_ph:
            self.fill_text(subtitle_ph, str(data.get("subtitle", "")))
        body_ph = self.get_placeholder(slide, idx=2)
        description = str(data.get("description") or "")
        if body_ph is not None:
            if description:
                self.fill_text(body_ph, description)
            elif data.get("bullets"):
                b = data["bullets"]
                if isinstance(b, str):
                    b = [p.strip() for p in re.split(r"[\n;]", b) if p.strip()]
                self.fill_bullets(body_ph, [str(x) for x in b])
        elif description:
            self.add_textbox(
                slide,
                description,
                left=0.65,
                top=4.5,
                width=12.0,
                height=1.5,
                font_size=18,
                color=self.theme["text_light"],
            )

    def add_agenda_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("agenda", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        item_bg_color = (
            self.theme["text_light"] if prefer_dark else self.theme["light_bg"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Agenda")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        items = data.get("items", [])
        if isinstance(items, str):
            items = [p.strip() for p in re.split(r"[\n;]", items) if p.strip()]
        start_y = 2.0
        for i, item in enumerate(items, 1):
            y_pos = start_y + (i - 1) * 0.9
            hex_shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON,
                Inches(0.75),
                Inches(y_pos),
                Inches(0.6),
                Inches(0.6),
            )
            hex_shape.fill.solid()
            hex_shape.fill.fore_color.rgb = _hex_to_rgb(self.theme["accent"])
            hex_shape.line.fill.background()
            tf = hex_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i)
            p.font.name = self.font_name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(self.theme["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.5),
                Inches(y_pos),
                Inches(8),
                Inches(0.6),
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = _hex_to_rgb(item_bg_color)
            item_bg.line.fill.background()
            self.add_textbox(
                slide,
                str(item),
                1.7,
                y_pos + 0.1,
                7.5,
                0.5,
                font_size=20,
                color=self.theme["text_dark"],
            )

    def add_timeline_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("timeline", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        body_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        secondary_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_secondary"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Timeline")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        steps = data.get("steps", [])
        num_steps = len(steps)
        if num_steps == 0:
            return
        step_width = 10.5 / num_steps
        start_x = 1.4
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x + 0.3),
            Inches(3.1),
            Inches(step_width * num_steps - 0.6),
            Inches(0.05),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_to_rgb(self.theme["accent"])
        line.line.fill.background()
        for i, step in enumerate(steps):
            if not isinstance(step, dict):
                step = {"title": str(step), "description": ""}
            x_pos = start_x + (i * step_width)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + step_width / 2 - 0.35),
                Inches(2.75),
                Inches(0.7),
                Inches(0.7),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = _hex_to_rgb(self.theme["accent"])
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.name = self.font_name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(self.theme["text_light"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE
            self.add_textbox(
                slide,
                str(step.get("title", f"Step {i + 1}")),
                x_pos,
                3.7,
                step_width,
                0.6,
                font_size=16,
                bold=True,
                color=body_color,
                alignment=PP_ALIGN.CENTER,
            )
            if step.get("description"):
                self.add_textbox(
                    slide,
                    str(step["description"]),
                    x_pos,
                    4.4,
                    step_width,
                    1.5,
                    font_size=12,
                    color=secondary_color,
                    alignment=PP_ALIGN.CENTER,
                )

    def add_icon_grid_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("icon-grid", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        body_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        secondary_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_secondary"]
        )
        circle_fill = (
            self.theme["text_light"] if prefer_dark else self.theme["light_bg"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Features")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        items = data.get("items", data.get("features", []))
        num_items = len(items)
        if num_items == 0:
            return
        if num_items <= 3:
            cols, rows = num_items, 1
        elif num_items <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 4, 2
        cell_width = 11 / cols
        cell_height = 2.2
        start_x = 1.2
        start_y = 1.8
        for i, item in enumerate(items[:8]):
            if not isinstance(item, dict):
                item = {"title": str(item), "description": "", "icon": ""}
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.5))
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + cell_width / 2 - 0.5),
                Inches(y_pos),
                Inches(1),
                Inches(1),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = _hex_to_rgb(circle_fill)
            circle.line.color.rgb = _hex_to_rgb(self.theme["accent"])
            circle.line.width = Pt(3)
            icon_raw = item.get("icon")
            if icon_raw:
                icon_text = icon_raw if len(str(icon_raw)) <= 2 else str(icon_raw)[0]
            else:
                icon_text = str(item.get("title", "?") or "?")[0].upper()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = icon_text
            p.font.name = self.font_name
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(self.theme["accent"])
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE
            self.add_textbox(
                slide,
                str(item.get("title", "")),
                x_pos,
                y_pos + 1.1,
                cell_width,
                0.5,
                font_size=14,
                bold=True,
                color=body_color,
                alignment=PP_ALIGN.CENTER,
            )
            if item.get("description"):
                self.add_textbox(
                    slide,
                    str(item["description"]),
                    x_pos,
                    y_pos + 1.55,
                    cell_width,
                    0.8,
                    font_size=11,
                    color=secondary_color,
                    alignment=PP_ALIGN.CENTER,
                )

    def add_stat_row_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("stat-row", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        label_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Key Metrics")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        stats = data.get("stats", [])
        num_stats = len(stats)
        if num_stats == 0:
            return
        stat_width = 11.5 / num_stats
        start_x = 0.9
        for i, stat in enumerate(stats):
            if not isinstance(stat, dict):
                stat = {"value": str(stat), "label": ""}
            x_pos = start_x + (i * stat_width)
            self.add_textbox(
                slide,
                str(stat.get("value", "0")),
                x_pos,
                2.5,
                stat_width - 0.3,
                1.5,
                font_size=56,
                bold=True,
                color=self.theme["accent"],
                alignment=PP_ALIGN.CENTER,
            )
            self.add_textbox(
                slide,
                str(stat.get("label", "")),
                x_pos,
                4.2,
                stat_width - 0.3,
                1.0,
                font_size=16,
                bold=True,
                color=label_color,
                alignment=PP_ALIGN.CENTER,
            )
            if i < num_stats - 1:
                divider = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos + stat_width - 0.15),
                    Inches(2.7),
                    Inches(0.02),
                    Inches(2.5),
                )
                divider.fill.solid()
                divider.fill.fore_color.rgb = _hex_to_rgb(self.theme["divider"])
                divider.line.fill.background()

    def add_pros_cons_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("pros-cons", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        body_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Pros & Cons")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        self.add_textbox(
            slide,
            str(data.get("pros_header", "Pros")),
            0.75,
            1.6,
            5.5,
            0.5,
            font_size=20,
            bold=True,
            color=self.theme["green"],
        )
        pros = data.get("pros", [])
        for i, pro in enumerate(pros):
            self.add_textbox(
                slide,
                f"✓  {pro}",
                0.75,
                2.2 + (i * 0.6),
                5.5,
                0.5,
                font_size=16,
                color=body_color,
            )
        self.add_textbox(
            slide,
            str(data.get("cons_header", "Cons")),
            7.0,
            1.6,
            5.5,
            0.5,
            font_size=20,
            bold=True,
            color=self.theme["red"],
        )
        cons = data.get("cons", [])
        for i, con in enumerate(cons):
            self.add_textbox(
                slide,
                f"✗  {con}",
                7.0,
                2.2 + (i * 0.6),
                5.5,
                0.5,
                font_size=16,
                color=body_color,
            )

    def add_comparison_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("comparison", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        label_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        diamond_bg = self.theme["dark_bg"] if prefer_dark else self.theme["accent"]
        self.add_textbox(
            slide,
            str(data.get("title", "Comparison")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(6.166),
            Inches(3.25),
            Inches(1),
            Inches(1),
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = _hex_to_rgb(diamond_bg)
        diamond.line.fill.background()
        tf = diamond.text_frame
        p = tf.paragraphs[0]
        p.text = "vs."
        p.font.name = self.font_name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(self.theme["text_light"])
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE
        self.add_textbox(
            slide,
            str(data.get("left_label", "Option A")),
            1.5,
            5.0,
            4.0,
            0.6,
            font_size=20,
            bold=True,
            color=label_color,
            alignment=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            slide,
            str(data.get("right_label", "Option B")),
            7.833,
            5.0,
            4.0,
            0.6,
            font_size=20,
            bold=True,
            color=label_color,
            alignment=PP_ALIGN.CENTER,
        )

    def add_checklist_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("checklist", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        body_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        unchecked_fill = (
            self.theme["text_light"] if prefer_dark else self.theme["light_bg"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Checklist")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        items = data.get("items", [])
        start_y = 1.8
        for i, item in enumerate(items):
            y_pos = start_y + (i * 0.7)
            checkbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.9),
                Inches(y_pos),
                Inches(0.35),
                Inches(0.35),
            )
            checkbox.fill.solid()
            is_checked = (
                item.get("checked", False) if isinstance(item, dict) else False
            )
            item_text = item.get("text", item) if isinstance(item, dict) else item
            if is_checked:
                checkbox.fill.fore_color.rgb = _hex_to_rgb(
                    self.theme["accent"]
                )
                tf = checkbox.text_frame
                p = tf.paragraphs[0]
                p.text = "✓"
                p.font.name = self.font_name
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = _hex_to_rgb(self.theme["text_light"])
                p.alignment = PP_ALIGN.CENTER
                tf.anchor = MSO_ANCHOR.MIDDLE
            else:
                checkbox.fill.fore_color.rgb = _hex_to_rgb(unchecked_fill)
            checkbox.line.color.rgb = _hex_to_rgb(self.theme["accent"])
            checkbox.line.width = Pt(2)
            self.add_textbox(
                slide,
                str(item_text),
                1.5,
                y_pos,
                10.0,
                0.4,
                font_size=16,
                color=body_color,
            )

    def add_logos_slide(self, data: dict[str, Any], *, prefer_dark: bool = False) -> None:
        slide = self._create_slide("logos", data, prefer_dark=prefer_dark)
        title_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_dark"]
        )
        subtitle_color = (
            self.theme["text_light"] if prefer_dark else self.theme["text_secondary"]
        )
        self.add_textbox(
            slide,
            str(data.get("title", "Our Partners")),
            0.83,
            0.59,
            10.0,
            0.8,
            font_size=36,
            bold=True,
            color=title_color,
        )
        if data.get("subtitle"):
            self.add_textbox(
                slide,
                str(data["subtitle"]),
                0.75,
                1.3,
                11.0,
                0.5,
                font_size=16,
                color=subtitle_color,
                alignment=PP_ALIGN.CENTER,
            )
        logos = data.get("logos", [])
        num_logos = len(logos)
        if num_logos <= 4:
            cols, rows = (num_logos or 1), 1
        elif num_logos <= 8:
            cols, rows = 4, 2
        else:
            cols, rows = 5, 2
        cell_width = 10 / max(cols, 1)
        cell_height = 1.6
        start_x = 1.7
        start_y = 2.5
        box_fill = self.theme["text_light"] if prefer_dark else self.theme["light_bg"]
        box_text = self.theme["text_secondary"]
        for i, logo in enumerate(logos[:10]):
            col = i % cols
            row = i // cols
            x_pos = start_x + (col * cell_width)
            y_pos = start_y + (row * (cell_height + 0.3))
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos),
                Inches(y_pos),
                Inches(cell_width - 0.4),
                Inches(cell_height - 0.3),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = _hex_to_rgb(box_fill)
            box.line.color.rgb = _hex_to_rgb(self.theme["divider"])
            box.line.width = Pt(1)
            logo_name = logo if isinstance(logo, str) else logo.get("name", "Company")
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = str(logo_name)
            p.font.name = self.font_name
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(box_text)
            p.alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE
