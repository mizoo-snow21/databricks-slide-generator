"""Import an existing .pptx file: extract per-slide text, build a Deck."""

from __future__ import annotations

from html import escape
from io import BytesIO
from typing import Any

from pptx import Presentation


def extract_slides_from_pptx(pptx_bytes: bytes) -> list[dict[str, Any]]:
    """Open a .pptx and return a list of slide dicts:
    [{layout, title, body_paragraphs, notes}, ...].

    Layout is inferred heuristically: first slide title, last closing, else content.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    out: list[dict[str, Any]] = []
    n_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        title = ""
        body_paras: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = (tf.text or "").strip()
            if not text:
                continue
            top = getattr(shape, "top", None)
            is_high = top is None or top < 1_500_000
            # First text frame with content becomes title if it's short and
            # appears near the top
            if not title and len(text) < 100 and is_high:
                title = text.split("\n")[0]
                for line in text.split("\n")[1:]:
                    line = line.strip()
                    if line:
                        body_paras.append(line)
            else:
                for line in text.split("\n"):
                    line = line.strip()
                    if line:
                        body_paras.append(line)
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        layout = "title" if i == 0 else ("closing" if i == n_slides - 1 else "content")
        out.append(
            {
                "layout": layout,
                "title": title or f"Slide {i + 1}",
                "body_paragraphs": body_paras,
                "notes": notes,
            }
        )
    return out


def build_html_from_extracted(
    slides: list[dict[str, Any]],
    deck_title: str = "Imported Deck",
) -> str:
    """Convert extracted slides into a deck HTML doc that goes through the
    standard sanitize → normalize → brand-inject pipeline.

    Each section is a basic light-content layout — the user is expected to
    refine it via the regenerate-slide / comment+apply flows.
    """
    sections: list[str] = []
    for i, s in enumerate(slides):
        sid = f"slide-{i + 1:02d}"
        oid_root = f"el-{sid}"
        paras = s.get("body_paragraphs") or []
        if not isinstance(paras, list):
            paras = []
        body_html = "".join(
            f'<p data-osd-id="{oid_root}-p{j}">{escape(str(p))}</p>'
            for j, p in enumerate(paras[:8])
        )
        raw_title = str(s.get("title") or f"Slide {i + 1}")
        layout = str(s.get("layout") or "content")
        title_html = (
            f'<h1 data-osd-id="{oid_root}-title">{escape(raw_title)}</h1>'
            if layout == "title"
            else f'<h2 data-osd-id="{oid_root}-title">{escape(raw_title)}</h2>'
        )
        sections.append(
            f'<section class="slide" '
            f'data-slide-id="{sid}" '
            f'data-osd-id="{oid_root}" '
            f'data-layout="{escape(layout)}">\n'
            f"{title_html}\n"
            f"{body_html}\n"
            "</section>"
        )
    body = "\n".join(sections)
    return (
        "<!doctype html>"
        "<html><head>"
        f"<title>{escape(deck_title)}</title>"
        "<style>"
        ":root { --osd-bg: #f5f3f0; --osd-text: #1b3139; --osd-accent: #ff3621; "
        "--osd-muted: #6b7280; "
        "--osd-font-display: 'DM Sans', sans-serif; --osd-font-body: 'DM Sans', sans-serif; }"
        "body { background: var(--osd-bg); color: var(--osd-text); "
        "font-family: var(--osd-font-body); margin: 0; }"
        "section.slide { width: 1920px; height: 1080px; padding: 100px; "
        "box-sizing: border-box; position: relative; }"
        "h1 { font-size: 96px; line-height: 1.05; margin: 0 0 24px; }"
        "h2 { font-size: 56px; line-height: 1.1; margin: 0 0 24px; }"
        "p { font-size: 22px; line-height: 1.5; margin: 0 0 16px; }"
        "</style>"
        "</head><body>"
        f"{body}"
        "</body></html>"
    )
